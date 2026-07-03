import sys
import io
from pathlib import Path
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Total slides: {len(prs.slides)}\n")

for idx, slide in enumerate(prs.slides):
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text.strip().replace("\n", " ")
    
    # Count shapes and images
    images_count = 0
    for shape in slide.shapes:
        if shape.name.startswith("Picture") or hasattr(shape, 'image'):
            images_count += 1
            
    print(f"Slide {idx + 1:02d} | Title: {repr(title):<60} | Shapes: {len(slide.shapes):<3} | Pictures: {images_count}")
