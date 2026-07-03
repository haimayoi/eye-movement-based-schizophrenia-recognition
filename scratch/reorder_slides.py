import sys
from pptx import Presentation

# 1. Open the Presentation
prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Loaded presentation with {len(prs.slides)} slides.")

# 2. Define the target order (1-indexed based on original slide positions)
# Introduction: 1 (Title), 2 (TOC), 3 (Motivation), 4 (Problem), 6 (Contributions)
# Related Works: 5 (Related Works & Limitations)
# Data Construction: 9 (T1), 10 (2D Filtering), 11 (Fixation Pos), 12 (Pupil Norm), 13 (T2 Features), 14 (Feature Dist), 15 (Top 8 Gaze), 16 (T2B Delta), 17 (Delta Dist)
# Proposed Method: 7 (5-Tier Pipeline Overview), 18 (T3 Tabular), 21 (T4A GNN-CEFAM), 23 (T4B BiCA-HS), 26 (T5 Meta-learner)
# Experimental Results: 8 (Evaluation Protocol - perfect start for results), 28 (Main Results Table), 29 (ROC Curves & DCA), 30 (Calibration Diagram), 31 (Ablation), 32 (Stimulus Category), 33 (Stats Testing), 34 (Multi-seed Stability), 19 (T3 Confusion Matrix), 20 (OOF Decision Space), 22 (t-SNE Embeddings), 24 (Threshold Sensitivity), 25 (Per-category P(SZ)), 27 (2D Ensemble Space), 35 (SHAP), 36 (SHAP Waterfall), 37 (Attention Vis), 38 (Feature Profile), 39 (Ensemble Decision Space), 40 (Case Study)
# Conclusion & Future Work: 41 (Curriculum Summary), 42 (Limitations), 43 (Conclusion)
# References: 44 (References), 45 (Thank you)
new_order_1based = [
    1, 2, 3, 4, 6,  # 1. Introduction
    5,              # 2. Related Works
    9, 10, 11, 12, 13, 14, 15, 16, 17, # 3. Data Construction
    7, 18, 21, 23, 26, # 4. Proposed Method
    8, 28, 29, 30, 31, 32, 33, 34, 19, 20, 22, 24, 25, 27, 35, 36, 37, 38, 39, 40, # 5. Experimental Results
    41, 42, 43,     # 6. Conclusion & Future Work
    44, 45          # 7. References
]

# Verify lengths and unique elements
assert len(new_order_1based) == len(prs.slides), f"Length mismatch: {len(new_order_1based)} vs {len(prs.slides)}"
assert set(new_order_1based) == set(range(1, len(prs.slides) + 1)), "Some slides are duplicated or missing in the target order!"

# 3. Access internal XML slide ID list and reorder
sldIdLst = prs.slides._sldIdLst
slide_layout_list = list(sldIdLst)

# Remove all slides from the list
for s_id in slide_layout_list:
    sldIdLst.remove(s_id)

# Re-append in the target order
for idx in new_order_1based:
    sldIdLst.append(slide_layout_list[idx - 1])

# 4. Save the presentation
prs.save("DAV_slides_v2_fixed.pptx")
print("\n✅ Presentation successfully reordered and saved as DAV_slides_v2_fixed.pptx!")
