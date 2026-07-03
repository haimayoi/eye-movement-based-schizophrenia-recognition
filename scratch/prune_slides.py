import sys
import io
import re
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# 1. Load the presentation
prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Loaded presentation with {len(prs.slides)} slides.")

# 2. Define the indices of slides to KEEP (1-indexed based on current reordered slides)
keep_indices_1based = [
    1,   # Title
    3,   # Motivation
    4,   # Problem Statement & Dataset (EMS)
    5,   # Key Contributions
    6,   # Related Works & Limitations
    7,   # T1 Data Preprocessing
    10,  # Pupil Dynamics & Normalization
    11,  # T2 Feature Engineering
    14,  # T2B Delta Features
    16,  # 5-Tier Pipeline Overview
    17,  # T3 Tabular Models
    18,  # T4A GNN-CEFAM
    19,  # T4B BiCA-HS
    20,  # T5 Meta-Learner
    21,  # Evaluation Protocol
    22,  # Main Results Table
    23,  # ROC Curves & DCA
    25,  # Ablation Study
    31,  # t-SNE Dimensionality Reduction
    35,  # SHAP Analysis
    37,  # Attention Visualization
    40,  # Case Study: SZ vs. HC
    42,  # Limitations & Potential Concerns
    43,  # Conclusion & Future Work
    44,  # References
    45   # Thank you
]

# Ensure keep indices are valid
assert len(keep_indices_1based) == 26, "Expected exactly 26 slides to be kept."
assert all(1 <= idx <= len(prs.slides) for idx in keep_indices_1based), "Invalid slide index in keep list."

# 3. Filter slides in the XML structure
sldIdLst = prs.slides._sldIdLst
slide_layout_list = list(sldIdLst)

# Clear list
for s_id in slide_layout_list:
    sldIdLst.remove(s_id)

# Re-append only the kept slides
for idx in keep_indices_1based:
    sldIdLst.append(slide_layout_list[idx - 1])

print(f"Pruned presentation to {len(prs.slides)} slides.")

# 4. Update the page footers of the remaining slides to show "X / 26"
modified_footers = 0
for idx, slide in enumerate(prs.slides):
    target_page_str = f"{idx + 1} / 26"
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    txt = r.text.strip()
                    # Match any slide footer pattern like "digit / digit"
                    if re.match(r"^\d+\s*/\s*\d+$", txt):
                        old_txt = r.text
                        r.text = target_page_str
                        modified_footers += 1
                        if modified_footers <= 5 or idx == 25:
                            print(f"  Slide {idx + 1} | Updated footer: {repr(old_txt)} -> {repr(target_page_str)}")

# 5. Save the pruned presentation under a new name to avoid Windows file locks
output_filename = "DAV_slides_v2_pruned.pptx"
prs.save(output_filename)
print(f"\n✅ Presentation successfully saved as {output_filename} with {modified_footers} updated footers!")
