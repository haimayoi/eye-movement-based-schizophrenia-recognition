import sys
import io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    prs = Presentation("DAV_slides_v2_pruned.pptx")
    print(f"Total slides: {len(prs.slides)}")
    
    print("\n--- Slide Layouts ---")
    for idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"Slide {idx+1:02d} | Title: {repr(title):<50} | Layout: {slide.slide_layout.name}")
        
    print("\n--- Available Slide Layouts ---")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: {repr(layout.name)}")
        
    print("\n--- Sample Shapes (Slide 3) ---")
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            print(f"  Shape Name: {repr(shape.name)} | Text: {repr(txt)}")
            
except Exception as e:
    print(f"Error: {e}")
