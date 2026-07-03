import sys
import os
from pathlib import Path
import io

# Force UTF-8 encoding for standard output
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

path = Path("DAV_slides_v2_fixed.pptx")
if not path.exists():
    print(f"Error: {path} not found.")
    exit(1)

prs = Presentation(path)
print(f"Total slides: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    print(f"\n==================================================")
    print(f"SLIDE {idx + 1}")
    
    # Try to find a title shape
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
    print(f"Title: {repr(title)}")
    
    images_count = 0
    text_content = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                p_text = paragraph.text.strip()
                if p_text:
                    text_content.append(p_text)
        if shape.name.startswith("Picture") or hasattr(shape, 'image'):
            images_count += 1
            
    print(f"Pictures count: {images_count}")
    print("Text Content:")
    for text in text_content[:15]:
        print(f"  - {text}")
    if len(text_content) > 15:
        print(f"  ... ({len(text_content) - 15} more lines)")
