import sys
import io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    prs = Presentation("DAV_slides_v2_fixed.pptx")
    print(f"Total slides: {len(prs.slides)}")
    
    # 1. Print slide order and layout names
    print("\n--- Current Slide Order & Layouts ---")
    for idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        layout_name = slide.slide_layout.name
        print(f"Slide {idx+1:02d} | Title: {repr(title):<50} | Layout: {layout_name}")
        
    # 2. Print all layout names in the presentation template
    print("\n--- Available Slide Layouts in Template ---")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: Name = {repr(layout.name)}")
        
    # 3. Inspect shapes in Slide 5 as a sample for footers
    print("\n--- Slide 5 Shapes & Text (Sample for footers) ---")
    slide = prs.slides[4]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            print(f"  Shape Name: {repr(shape.name)} | Text: {repr(text)}")
            
except Exception as e:
    print(f"Error: {e}")
