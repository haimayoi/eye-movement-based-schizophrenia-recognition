import sys
import io
import re
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# 1. Load the presentation
prs = Presentation("DAV_slides_v2_pruned.pptx")
total_slides = len(prs.slides)
print(f"Loaded presentation with {total_slides} slides.")

# 2. Get Layout index 5 ('1_Comparison')
layout_comparison = prs.slide_layouts[5]
print(f"Using slide layout: {layout_comparison.name}")

# 3. Process each slide
deleted_footers_count = 0
updated_pages_count = 0

for idx, slide in enumerate(prs.slides):
    # Set slide layout to 1_Comparison
    slide.slide_layout = layout_comparison
    
    # Identify shapes to delete and slide number shape
    shapes_to_delete = []
    slide_number_shape = None
    
    target_page_str = f"{idx + 1} / {total_slides}"
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Aggregate paragraph text
            text_content = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            
            # Check for footer patterns: e.g., "Thesis Defense", "Defense Academy", etc.
            if any(term in text_content for term in ["Thesis Defense", "Defense Academy", "University of", "HUST", "Faculty of"]):
                shapes_to_delete.append(shape)
            # Check if this is the slide page number shape
            elif re.match(r"^\d+\s*/\s*\d+$", text_content.strip()):
                slide_number_shape = shape
                
    # Delete the footer shapes
    for shape in shapes_to_delete:
        try:
            slide.shapes._spTree.remove(shape._element)
            deleted_footers_count += 1
        except Exception as e:
            print(f"  Error deleting shape {shape.name}: {e}")
            
    # Update page number
    if slide_number_shape:
        # Clear existing text and set to target page number
        slide_number_shape.text_frame.text = target_page_str
        # Formatting
        p = slide_number_shape.text_frame.paragraphs[0]
        if len(p.runs) > 0:
            p.runs[0].font.name = "Arial"
        updated_pages_count += 1
    else:
        # If no page number shape exists on this slide (like Slide 1 originally), 
        # let's check if we should add one or search for shapes that look like slide numbers
        pass

# 4. Save the presentation
output_file = "DAV_slides_v2_pruned.pptx"
prs.save(output_file)

print(f"\nCompleted processing!")
print(f"- Set all layouts to '1_Comparison'")
print(f"- Deleted {deleted_footers_count} footer shapes.")
print(f"- Updated {updated_pages_count} slide page numbers to match f'X / {total_slides}' format.")
print(f"- Saved presentation to {output_file}")
