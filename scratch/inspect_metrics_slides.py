import sys
import io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Slide width: {prs.slide_width / 9144000:.2f} inches")
print(f"Slide height: {prs.slide_height / 9144000:.2f} inches")

for idx in [42, 43]:
    slide = prs.slides[idx - 1]
    print(f"\n========================================= Slide {idx} =========================================")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for run_idx, r in enumerate(p.runs):
                    txt = r.text.strip()
                    if any(w in txt for w in ["0.9454", "6.68%", "p<0.05", "bootstrap"]):
                        print(f"Shape: {shape.name} | Paragraph {p_idx} | Run {run_idx}: {repr(r.text)}")
