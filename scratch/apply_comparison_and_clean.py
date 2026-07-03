import sys
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# 1. Load the presentation
prs = Presentation("DAV_slides_v2_pruned.pptx")
total_slides = len(prs.slides)
print(f"Loaded presentation with {total_slides} slides.")

# 2. Get Layout index 5 ('1_Comparison')
layout_comparison = prs.slide_layouts[5]
print(f"Target slide layout: {layout_comparison.name}")

deleted_footers_count = 0
updated_pages_count = 0
added_pages_count = 0

# 3. Process each slide
for idx, slide in enumerate(prs.slides):
    # Change slide layout relationship target to 1_Comparison
    rel_id = None
    for rId, rel in slide.part.rels.items():
        if "slideLayout" in rel.reltype:
            rel_id = rId
            break
            
    if rel_id:
        rel = slide.part.rels[rel_id]
        rel._target = layout_comparison.part
        # Clear lazyproperty cache so python-pptx serializes it correctly
        rel.__dict__.pop('target_ref', None)
        rel.__dict__.pop('target_partname', None)
        rel.__dict__.pop('target_part', None)
        
    # Standard target page number
    target_page_str = f"{idx + 1} / {total_slides}"
    
    # Identify shapes to delete or update
    shapes_to_delete = []
    slide_number_shape = None
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_content = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            
            # Check for footer text (Thesis, HUST, DAV, etc.)
            if any(term in text_content for term in ["Thesis Defense", "Defense Academy", "University of", "HUST", "Faculty of"]):
                shapes_to_delete.append(shape)
            # Check if this shape is a slide number
            elif re.match(r"^\d+\s*/\s*\d+$", text_content.strip()):
                slide_number_shape = shape
            # Or check if it is a bottom-right textbox containing digits (fallback search)
            elif shape.left is not None and shape.top is not None:
                if shape.left > Inches(8.0) and shape.top > Inches(6.5) and re.match(r"^\d+$", text_content.strip()):
                    slide_number_shape = shape
                
    # Delete the footer shapes
    for shape in shapes_to_delete:
        try:
            slide.shapes._spTree.remove(shape._element)
            deleted_footers_count += 1
        except Exception as e:
            print(f"  Error deleting shape {shape.name}: {e}")
            
    # Update or add page number
    if slide_number_shape:
        slide_number_shape.text_frame.text = target_page_str
        p = slide_number_shape.text_frame.paragraphs[0]
        p.alignment = 2  # Right aligned
        if len(p.runs) > 0:
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(128, 128, 128) # Gray color
        updated_pages_count += 1
    else:
        # Create a new standardized page number box at the bottom right
        left_coord = Inches(8.8)
        top_coord = Inches(7.0)
        width_coord = Inches(1.0)
        height_coord = Inches(0.4)
        
        txBox = slide.shapes.add_textbox(left_coord, top_coord, width_coord, height_coord)
        tf = txBox.text_frame
        tf.text = target_page_str
        
        p = tf.paragraphs[0]
        p.alignment = 2  # Right aligned
        if len(p.runs) > 0:
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.bold = False
            run.font.color.rgb = RGBColor(128, 128, 128) # Gray color
            
        added_pages_count += 1

# 4. Save presentation
prs.save("DAV_slides_v2_pruned.pptx")

print(f"\nProcessing Completed successfully!")
print(f"- Set all layouts to '1_Comparison' (HUST template design).")
print(f"- Deleted {deleted_footers_count} footer shapes.")
print(f"- Updated {updated_pages_count} existing slide numbers.")
print(f"- Added {added_pages_count} new slide numbers.")
print(f"- Saved cleaned slides to DAV_slides_v2_pruned.pptx")
