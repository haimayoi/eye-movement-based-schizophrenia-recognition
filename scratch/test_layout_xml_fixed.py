import sys
from pptx import Presentation

try:
    prs = Presentation("DAV_slides_v2_pruned.pptx")
    slide = prs.slides[0]
    print(f"Initial layout: {slide.slide_layout.name}")
    
    # Target layout: 1_Comparison (Index 5)
    target_layout = prs.slide_layouts[5]
    
    # Find layout relationship in slide's relationships
    rel_id = None
    for rId, rel in slide.part.rels.items():
        if "slideLayout" in rel.reltype:
            rel_id = rId
            print(f"Found layout relationship: rId={rId}, target_ref={rel.target_ref}")
            break
            
    if rel_id:
        rel = slide.part.rels[rel_id]
        # Change relationship target to target layout's part
        rel._target = target_layout.part
        # Clear the lazyproperty cache in relationship object dictionary
        rel.__dict__.pop('target_ref', None)
        rel.__dict__.pop('target_partname', None)
        rel.__dict__.pop('target_part', None)
        print(f"Cleared cache. New target_ref: {rel.target_ref}")
        
    prs.save("scratch/test_layout_out.pptx")
    print("Saved test presentation successfully!")
    
    # Reload and check
    prs2 = Presentation("scratch/test_layout_out.pptx")
    print(f"Reloaded slide 1 layout: {prs2.slides[0].slide_layout.name}")
    
except Exception as e:
    print(f"Error: {e}")
