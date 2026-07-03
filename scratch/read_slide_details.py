import sys
import io
from pathlib import Path
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

prs = Presentation("DAV_slides_v2_fixed.pptx")

slides_to_read = [28, 29, 30, 31, 32, 38, 39, 40]

for idx in slides_to_read:
    if idx > len(prs.slides):
        continue
    slide = prs.slides[idx - 1]
    print(f"\n==================================================")
    print(f"SLIDE {idx} | Title: {slide.shapes.title.text if slide.shapes.title else ''}")
    
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    text_content.append(p.text.strip())
                    
    print("Text:")
    for text in text_content:
        print(f"  - {text}")
