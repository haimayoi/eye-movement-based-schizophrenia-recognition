from pptx import Presentation

prs = Presentation("DAV_slides_v2_fixed.pptx")
slide = prs.slides[40]  # Slide 41 (0-indexed is 40)

print(f"Slide 41 title: {slide.shapes.title.text if slide.shapes.title else 'No Title'}")
print(f"Total shapes: {len(slide.shapes)}")

for i, shape in enumerate(slide.shapes):
    print(f"\nShape {i+1}: Name={shape.name}, Type={shape.shape_type}")
    if shape.has_text_frame:
        print(f"  Text Frame Paragraphs: {len(shape.text_frame.paragraphs)}")
        for j, p in enumerate(shape.text_frame.paragraphs):
            print(f"    P {j+1}: {repr(p.text)}")
