import sys
import io
import re
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Loaded presentation with {len(prs.slides)} slides.")

modified_count = 0
for idx, slide in enumerate(prs.slides):
    target_page_str = f"{idx + 1} / 45"
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    txt = r.text.strip()
                    # Match any string of the form "digits / 45" or "digits  / 45"
                    if re.match(r"^\d+\s*/\s*45$", txt):
                        if txt != target_page_str:
                            old_txt = r.text
                            r.text = target_page_str
                            modified_count += 1
                            print(f"Slide {idx + 1} | Updated footer: {repr(old_txt)} -> {repr(target_page_str)}")

prs.save("DAV_slides_v2_fixed.pptx")
print(f"\nSuccessfully corrected {modified_count} slide page number footers!")
