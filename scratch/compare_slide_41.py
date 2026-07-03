import sys
import io
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

for name in ["DAV_slides.pptx", "DAV_slides_v2.pptx"]:
    try:
        prs = Presentation(name)
        print(f"\n==================================================")
        print(f"File: {name} | Total Slides: {len(prs.slides)}")
        # Let's search for "DAV CURRICULUM" or similar in all slides of this file
        for idx, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else ""
            if "CURRICULUM" in title.upper() or "COVERAGE" in title.upper():
                print(f"Found 'Curriculum' slide at index {idx+1}!")
                print(f"Slide {idx+1} title: {title}")
                print(f"Total shapes: {len(slide.shapes)}")
                for i, shape in enumerate(slide.shapes):
                    if shape.has_text_frame:
                        txt = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                        if txt:
                            print(f"  Shape {i+1}: {repr(txt)}")
    except Exception as e:
        print(f"Error reading {name}: {e}")
