import sys
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# 1. Open the Presentation
prs = Presentation("DAV_slides_v2_fixed.pptx")
print(f"Loaded presentation with {len(prs.slides)} slides.")

# 2. Key replacements for outdated metrics and overclaims
replacements = {
    "0.9454": "0.9549",
    "0.0099": "0.0039",
    "+6.68%": "+6.95%",
    "surpass SOTA with statistical significance (bootstrap CI, p<0.05)": 
        "exceed the reported state-of-the-art MSNet (reported test AUC=0.8854)",
    "outperform the state-of-the-art MSNet":
        "exceed the reported state-of-the-art MSNet"
}

modified_runs = 0
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    orig_text = r.text
                    new_text = orig_text
                    for old, new in replacements.items():
                        if old in new_text:
                            new_text = new_text.replace(old, new)
                    if new_text != orig_text:
                        r.text = new_text
                        modified_runs += 1
                        print(f"Slide {idx+1} | Replaced text: {repr(orig_text)} -> {repr(new_text)}")

print(f"\nCompleted {modified_runs} text replacements.")

# 3. Populate Slide 41 (DAV CURRICULUM COVERAGE SUMMARY)
# Index 40 is Slide 41
slide_41 = prs.slides[40]

# Remove the empty placeholder (Shape 2: Content Placeholder 2)
shape_to_remove = None
for shape in slide_41.shapes:
    if shape.name == "Content Placeholder 2":
        shape_to_remove = shape
        break

if shape_to_remove:
    print("\nRemoving empty content placeholder on Slide 41...")
    # Safe XML element removal
    spTree = slide_41.shapes._spTree
    spTree.remove(shape_to_remove._element)
    print("Empty placeholder removed.")

# Define table coordinates (centered, fits standard 4:3 slide)
left = Inches(1.0)
top = Inches(2.0)
width = Inches(8.0)
height = Inches(4.5)

rows = 6
cols = 3

print("Adding curriculum coverage summary table on Slide 41...")
table_shape = slide_41.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(2.3)  # Requirement
table.columns[1].width = Inches(4.2)  # Focus / Techniques
table.columns[2].width = Inches(1.5)  # Slides

# Table Data
table_data = [
    ["DAV Requirement", "Focus / Deep Learning Techniques", "Slide(s)"],
    ["C1: Evaluation & Thresholds", "ST-GNN vs. BiCA-HS comparisons, decision space, probability calibration, ECE/Brier metric tracking.", "24, 25, 27, 30"],
    ["C2: Preprocessing & Data", "Spatial screen-boundary filters, temporal 50ms cutoff, pupil size subject-wise normalization.", "10, 11, 12, 14"],
    ["C3: Feature Engineering", "Top 8 discriminative features, contextual category-mean features, cross-category contrastive delta features.", "15, 16, 17"],
    ["C5/C6: Dimension & Space", "t-SNE dimensionality reduction of visual node features, OOF single-model decision space mapping.", "20, 22"],
    ["C7: Interpretability & Case", "SHAP beeswarm importance, SHAP waterfall individual explanation, subject-level case study comparisons (SZ vs. HC).", "36, 38, 39, 40"]
]

# Write data and format text
for row_idx, row in enumerate(table_data):
    for col_idx, text in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        
        # Style formatting
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Arial"
        
        if row_idx == 0:
            # Header Row
            run.font.size = Pt(12)
            run.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(211, 47, 47)  # Dark Red HUST-like color
            run.font.color.rgb = RGBColor(255, 255, 255)      # White text
        else:
            # Data Rows
            run.font.size = Pt(10)
            if col_idx == 0:
                run.font.bold = True
            cell.fill.solid()
            if row_idx % 2 == 1:
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245) # Light gray zebra
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # White

print("Slide 41 table added and formatted successfully.")

# Save presentation
prs.save("DAV_slides_v2_fixed.pptx")
print("\n✅ Presentation saved as DAV_slides_v2_fixed.pptx!")
