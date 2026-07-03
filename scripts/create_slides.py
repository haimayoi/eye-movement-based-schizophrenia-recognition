# -*- coding: utf-8 -*-
"""
Create DAV_slides.pptx for thesis defense
"Eye Movement-Based Schizophrenia Recognition using Hierarchical Multi-Stream Learning"
English, HUST template
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

BASE = r"D:\DAV\Eye Movement-Based Schizophrenia Recognition"
TEMPLATE = os.path.join(BASE, "HUST_PPT_template_2022_RED_4x3-3.pptx")
OUTPUT = os.path.join(BASE, "DAV_slides.pptx")

# HUST colors
RED = RGBColor(0xC0, 0x00, 0x00)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BLUE_ACCENT = RGBColor(0x1F, 0x49, 0x7D)
GREEN = RGBColor(0x00, 0x70, 0x00)

# Slide dimensions (EMU): 10" x 7.5"
SW = 9144000
SH = 6858000

def img(rel_path):
    """Return absolute path to image, or None if not found."""
    p = os.path.join(BASE, rel_path)
    return p if os.path.exists(p) else None

def delete_all_slides(prs):
    """Remove all existing slides from the presentation."""
    xml_slides = prs.slides._sldIdLst
    # Remove from end to avoid index issues
    slide_ids = list(xml_slides)
    for sId in slide_ids:
        rId = sId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sId)

def add_slide(prs, layout_index):
    """Add a new slide using the given layout index."""
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    return slide

def set_title(slide, text, font_size=Pt(28), bold=True, color=None):
    """Set the title placeholder text."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = font_size
                    run.font.bold = bold
                    if color:
                        run.font.color.rgb = color
            return ph
    return None

def set_content(slide, text_or_bullets, font_size=Pt(18), idx=13):
    """Set content placeholder (idx=13) with text or bullet list."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.word_wrap = True
            tf.clear()
            if isinstance(text_or_bullets, list):
                for i, bullet in enumerate(text_or_bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    if bullet.startswith('  '):
                        p.level = 1
                        p.text = bullet.strip()
                    for run in p.runs:
                        run.font.size = font_size
            else:
                tf.paragraphs[0].text = text_or_bullets
                for run in tf.paragraphs[0].runs:
                    run.font.size = font_size
            return ph
    return None

def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(16), bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, bg_color=None, word_wrap=True):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox

def add_picture_safe(slide, rel_path, left, top, width, height):
    """Add a picture if it exists, return shape or None."""
    p = img(rel_path)
    if p:
        try:
            return slide.shapes.add_picture(p, left, top, width, height)
        except Exception as e:
            print("WARNING: Could not add image %s: %s" % (rel_path, e))
    else:
        print("WARNING: Missing image: %s" % rel_path)
    return None

def add_rect(slide, left, top, width, height, fill_color=RED,
             line_color=None, text=None, font_size=Pt(12),
             text_color=WHITE, bold=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape

def add_table_slide(slide, headers, rows, left, top, width, height,
                    header_bg=RED, header_fg=WHITE,
                    row_bg=WHITE, alt_bg=LIGHT_GRAY,
                    font_size=Pt(13)):
    """Add a formatted table to a slide."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Style header
    for col_idx, hdr in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = font_size
                run.font.bold = True
                run.font.color.rgb = header_fg

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        bg = alt_bg if row_idx % 2 == 1 else row_bg
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = font_size

    return table

def add_slide_number(slide, slide_num, total=27):
    """Add slide number to bottom right."""
    add_textbox(
        slide, "%d / %d" % (slide_num, total),
        left=SW - Inches(1.2), top=SH - Inches(0.35),
        width=Inches(1.0), height=Inches(0.3),
        font_size=Pt(10), color=DARK_GRAY,
        align=PP_ALIGN.RIGHT
    )

def add_footer_line(slide, text="Thesis Defense — DAV / HUST 2025"):
    """Add a thin footer line."""
    add_textbox(
        slide, text,
        left=Inches(0.2), top=SH - Inches(0.35),
        width=Inches(7), height=Inches(0.3),
        font_size=Pt(9), color=DARK_GRAY,
        align=PP_ALIGN.LEFT
    )

# ============================================================
# MAIN
# ============================================================
print("Loading template...")
prs = Presentation(TEMPLATE)
print("Template loaded. Slides:", len(prs.slides))
print("Deleting example slides...")
delete_all_slides(prs)
print("Slides after delete:", len(prs.slides))

TOTAL_SLIDES = 27

# ============================================================
# SLIDE 1: Trang bìa (Title Slide)
# ============================================================
slide = add_slide(prs, 1)  # 1_Title Slide

# Title placeholder (idx=0) is around the center-right area
# Add custom textboxes for full cover design
add_rect(slide, 0, 0, SW, Inches(0.9), fill_color=RED)
add_textbox(
    slide,
    "DEFENSE ACADEMY OF VIETNAM (DAV)",
    left=Inches(0.3), top=Inches(0.05),
    width=Inches(9.4), height=Inches(0.5),
    font_size=Pt(16), bold=True, color=WHITE,
    align=PP_ALIGN.CENTER
)
add_textbox(
    slide,
    "FACULTY OF INFORMATION TECHNOLOGY",
    left=Inches(0.3), top=Inches(0.45),
    width=Inches(9.4), height=Inches(0.4),
    font_size=Pt(13), bold=False, color=WHITE,
    align=PP_ALIGN.CENTER
)

# Main title box
add_rect(slide, Inches(0.5), Inches(1.1), SW - Inches(1.0), Inches(2.2),
         fill_color=RGBColor(0xF5, 0xE6, 0xE6))
add_textbox(
    slide,
    "EYE MOVEMENT-BASED SCHIZOPHRENIA RECOGNITION\nUSING HIERARCHICAL MULTI-STREAM LEARNING",
    left=Inches(0.6), top=Inches(1.15),
    width=SW - Inches(1.2), height=Inches(2.1),
    font_size=Pt(26), bold=True, color=RED,
    align=PP_ALIGN.CENTER
)

# Subtitle
add_textbox(
    slide,
    "A 5-Tier Pipeline: Preprocessing → Feature Engineering → Tabular → Deep Learning → Meta-Learner",
    left=Inches(0.5), top=Inches(3.35),
    width=SW - Inches(1.0), height=Inches(0.8),
    font_size=Pt(16), bold=False, color=BLUE_ACCENT,
    align=PP_ALIGN.CENTER
)

# Author info
add_textbox(
    slide,
    "Student: [Student Name]\nSupervisor: [Supervisor Name]\nDefense Academy of Vietnam (DAV) — 2025",
    left=Inches(2.0), top=Inches(4.3),
    width=Inches(6.0), height=Inches(1.2),
    font_size=Pt(16), bold=False, color=DARK_GRAY,
    align=PP_ALIGN.CENTER
)

# Bottom bar
add_rect(slide, 0, SH - Inches(0.6), SW, Inches(0.6), fill_color=RED)
add_textbox(
    slide, "Hanoi, 2025",
    left=Inches(0.3), top=SH - Inches(0.55),
    width=Inches(9.4), height=Inches(0.5),
    font_size=Pt(13), bold=False, color=WHITE,
    align=PP_ALIGN.CENTER
)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: Mục lục (Agenda)
# ============================================================
slide = add_slide(prs, 3)  # 1_Section Header

set_title(slide, "TABLE OF CONTENTS", font_size=Pt(28), bold=True)
bullets = [
    "1.  Motivation & Problem Statement",
    "2.  Dataset (EMS) & Research Questions",
    "3.  Related Works & Research Gap",
    "4.  Key Contributions",
    "5.  5-Tier Hierarchical Pipeline",
    "6.  Experiments & Results",
    "7.  Explainability Analysis",
    "8.  Conclusion & Future Work",
]
set_content(slide, bullets, font_size=Pt(20), idx=13)
add_footer_line(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3: Động lực nghiên cứu
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "MOTIVATION", font_size=Pt(26))
bullets = [
    "Schizophrenia (SZ) affects ~1% of the global population (~75 million people)",
    "Current clinical diagnosis is entirely subjective (interview-based, DSM-5) — prone to inconsistency",
    "Eye-tracking is a non-invasive, objective biomarker: SZ patients show clear oculomotor abnormalities",
    "Challenge: No existing model integrates multi-level features (handcrafted + graph + raw sequences) jointly",
    "Opportunity: EMS dataset (Song et al. 2025) — 208 subjects, 4 stimulus categories, largest published",
    "Goal: Build a 5-tier hierarchical pipeline that fully exploits eye-tracking signals for SZ/HC classification",
]
set_content(slide, bullets, font_size=Pt(17), idx=13)
add_footer_line(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4: Bài toán & Dataset
# ============================================================
slide = add_slide(prs, 4)  # 1_Two Content

set_title(slide, "PROBLEM STATEMENT & DATASET (EMS)", font_size=Pt(26))

# Left: Dataset info
left_bullets = [
    "EMS Dataset (Song et al. 2025)",
    "  208 subjects: 160 train + 48 test",
    "  80 SZ + 80 HC (training labels)",
    "  48 test: labels withheld (competition)",
    "  281,037 total fixation records",
    "",
    "4 stimulus categories:",
    "  Social — portrait images of people",
    "  Natural — landscape / nature scenes",
    "  Synthetic — texture / synthetic images",
    "  Manipulated — edited / altered images",
    "",
    "~100 stimuli per subject",
    "Fixation sequences: up to 200 fixations/trial",
]
# Set left placeholder (idx=1)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, b in enumerate(left_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(15)
        break

# Right: columns info as text + image
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 2:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        right_bullets = [
            "Raw data columns:",
            "  IMAGE, FIX_INDEX, FIX_DURATION",
            "  FIX_X, FIX_Y, FIX_PUPIL",
            "  Subject_ID, Label, Is_Test",
            "",
            "Evaluation protocol:",
            "  4-fold GroupKFold",
            "  (from official EMS Train_Valid.xlsx)",
            "  3 seeds: 42, 123, 456",
            "  OOF (Out-Of-Fold) predictions",
            "",
            "Reference SOTA:",
            "  MSNet: AUC=0.8854, ACC=81.25%",
        ]
        for i, b in enumerate(right_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(15)
        break

add_footer_line(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5: Các nghiên cứu liên quan
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "RELATED WORKS & LIMITATIONS", font_size=Pt(24))

# Add a comparison table
headers = ["Method", "Features", "Architecture", "Limitation"]
rows = [
    ["Itti et al. (2000)", "Handcrafted", "Saliency map", "Cannot learn from data"],
    ["Nakashima et al.", "CNN visual", "ResNet", "Ignores temporal sequence"],
    ["Li et al. (GNN)", "Graph nodes", "GCN", "Lacks integrated biomarkers"],
    ["MSNet (SOTA)", "Multi-scale", "CNN+Fusion", "AUC=0.8854, does not leverage OT"],
    ["Ours (5-tier)", "135 features\n+graph+seq", "CEFAM+BiCA-HS", "Most comprehensive pipeline"],
]
add_table_slide(
    slide, headers, rows,
    left=Inches(0.3), top=Inches(1.0),
    width=Inches(9.4), height=Inches(4.8),
    font_size=Pt(14)
)

# Highlight our row
add_textbox(
    slide,
    "Research gap: No prior study simultaneously combines handcrafted features, GNN spatiotemporal graph, and Transformer sequence learning on the EMS dataset",
    left=Inches(0.3), top=Inches(5.85),
    width=Inches(9.4), height=Inches(0.5),
    font_size=Pt(13), bold=True, color=RED,
    align=PP_ALIGN.LEFT
)

add_footer_line(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6: Đóng góp chính
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "KEY CONTRIBUTIONS", font_size=Pt(28))
bullets = [
    "1. 5-Tier Hierarchical Pipeline: Preprocessing → Feature Engineering → Tabular → DL → Meta-Learner",
    "2. CEFAM (Cross-attention Expert-Fixation Attention Module): Bidirectional attention between expert knowledge stream and GNN fixation nodes",
    "3. BiCA-HS (Bidirectional Cross-Attention Hybrid Stream): Combines Transformer (raw sequences) + Biomarker MLP (15 stimulus-level features) via bidirectional cross-attention",
    "4. Delta Features: 60 contrast features across stimulus-category pairs — captures SZ-specific viewing anomalies by visual context",
    "5. Comprehensive Evaluation: Bootstrap CI, DeLong test, multi-seed stability (3 seeds), 4-component ablation study",
]
set_content(slide, bullets, font_size=Pt(17), idx=13)

add_footer_line(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7: Tổng quan Pipeline (diagram)
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "5-TIER PIPELINE OVERVIEW", font_size=Pt(26))

# Draw pipeline boxes
tiers = [
    ("TIER 1\nPreprocessing", "Spatial/Temporal\nFilter"),
    ("TIER 2\nFeature Eng.", "15 → 135\nFeatures"),
    ("TIER 3\nTabular", "XGBoost\nCatBoost\nLightGBM"),
    ("TIER 4\nDL Models", "GNN-CEFAM\nBiCA-HS"),
    ("TIER 5\nMeta-Learner", "Weighted Avg\n/ Log. Reg."),
]

box_w = Inches(1.6)
box_h = Inches(1.5)
gap = Inches(0.25)
start_x = Inches(0.3)
y_top = Inches(1.2)

colors_tier = [
    RGBColor(0x2E, 0x75, 0xB6),  # blue
    RGBColor(0x70, 0xAD, 0x47),  # green
    RGBColor(0xFF, 0x99, 0x00),  # orange
    RGBColor(0xC0, 0x00, 0x00),  # red
    RGBColor(0x7B, 0x2F, 0xBE),  # purple
]

for i, (title, subtitle) in enumerate(tiers):
    x = start_x + i * (box_w + gap)
    # Main box
    add_rect(slide, x, y_top, box_w, box_h,
             fill_color=colors_tier[i],
             text=title,
             font_size=Pt(13), bold=True, text_color=WHITE)
    # Subtitle box
    add_rect(slide, x, y_top + box_h,
             box_w, Inches(0.7),
             fill_color=LIGHT_GRAY,
             text=subtitle,
             font_size=Pt(11), bold=False, text_color=DARK_GRAY)
    # Arrow (except last)
    if i < len(tiers) - 1:
        arr_x = x + box_w + Inches(0.03)
        arr_y = y_top + box_h / 2 - Inches(0.1)
        add_textbox(slide, "→",
                    left=arr_x, top=arr_y,
                    width=Inches(0.2), height=Inches(0.3),
                    font_size=Pt(20), bold=True, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER)

# Data flow labels
add_textbox(
    slide,
    "clean_fixations.parquet    →    subject_features.csv    →    OOF predictions    →    Final Score",
    left=Inches(0.3), top=Inches(3.05),
    width=Inches(9.4), height=Inches(0.4),
    font_size=Pt(12), bold=False, color=DARK_GRAY,
    align=PP_ALIGN.CENTER
)

# Input/Output labels
inputs_outputs = [
    (start_x + Inches(0.2), "281K fixations\nraw EMS data"),
    (start_x + Inches(1.85) + gap, "135 features\n/subject"),
    (start_x + Inches(3.7) + gap*2, "OOF AUC\nper model"),
    (start_x + Inches(5.55) + gap*3, "Deep OOF\n(4-fold)"),
    (start_x + Inches(7.4) + gap*4, "Best ensemble\nprediction"),
]
for x_pos, label in inputs_outputs:
    add_textbox(slide, label,
                left=x_pos, top=Inches(3.5),
                width=box_w, height=Inches(0.7),
                font_size=Pt(10), color=BLUE_ACCENT,
                align=PP_ALIGN.CENTER)

# Key innovations
add_textbox(
    slide,
    "Innovation: CEFAM attention (Tier 4A) | BiCA-HS dual-stream (Tier 4B) | Delta features (Tier 2B) | OOF 4-fold evaluation",
    left=Inches(0.3), top=Inches(6.3),
    width=Inches(9.4), height=Inches(0.4),
    font_size=Pt(12), bold=True, color=RED,
    align=PP_ALIGN.CENTER
)

add_footer_line(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8: Tier 1 - Tiền xử lý
# ============================================================
slide = add_slide(prs, 4)  # Two Content

set_title(slide, "TIER 1 — DATA PREPROCESSING", font_size=Pt(26))

# Left: text
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Spatial Filter:",
            "  Remove fixations outside 1024×768 screen",
            "  → Retain valid fixation points",
            "",
            "Temporal Filter:",
            "  Remove too-short fixations (<50ms)",
            "  Remove too-long fixations (>2000ms)",
            "  → Retain naturalistic fixations",
            "",
            "Output:",
            "  Input: raw EMS fixation records",
            "  Output: clean_fixations.parquet",
            "  Total: 281,037 clean fixation records",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(15)
        break

# Right: image
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 2:
        ph_left = ph.left
        ph_top = ph.top
        ph_width = ph.width
        ph_height = ph.height
        break

add_picture_safe(slide, "results/tier1_2d_filtering.png",
                 Inches(5.0), Inches(1.3), Inches(4.7), Inches(4.5))

add_footer_line(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9: Tier 2 - Feature Engineering
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "TIER 2 — FEATURE ENGINEERING", font_size=Pt(26))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Tier 2A — 15 stimulus-level features:",
            "  fix_count, fix_dur_mean, fix_dur_median, fix_dur_std",
            "  sacc_amp_mean, sacc_amp_std, sa_fdr_mean, sacc_angle_mean",
            "  scanpath_length, convex_hull_area, center_bias",
            "  spatial_entropy, pupil_mean, pupil_std, pupil_cv",
            "",
            "Tier 2B — 120 subject-level features:",
            "  60 category-means: 4 categories × 15 features",
            "  60 delta features: 4 pairs × 15 features",
            "  (contrast between stimulus categories)",
            "",
            "Total: 15 + 120 = 135 features for Tabular/GNN",
            "BiCA-HS: uses only 15 stimulus-level features (d_bio=15)",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier2_correlation_heatmap.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(4.8))

add_footer_line(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10: Tier 2 - Delta Features
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "TIER 2B — DELTA FEATURES", font_size=Pt(26))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Insight: SZ patients respond differently to different image categories",
            "",
            "4 delta pairs (4 pairs × 15 = 60 features):",
            "  Social − Natural",
            "  Social − Synthetic",
            "  Social − Manipulated",
            "  Natural − Synthetic",
            "",
            "Example: Δ(fix_dur_mean) — do SZ fixate longer on social",
            "images compared to natural scenes?",
            "",
            "Helps model capture context-specific viewing anomalies",
            "No data leakage (fold-wise normalization applied)",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(15)
        break

add_picture_safe(slide, "results/tier2_delta_mechanism.png",
                 Inches(5.0), Inches(1.2), Inches(4.7), Inches(2.5))
add_picture_safe(slide, "results/tier2_delta_density.png",
                 Inches(5.0), Inches(3.8), Inches(4.7), Inches(2.5))

add_footer_line(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11: Tier 3 - Tabular Models
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "TIER 3 — TABULAR MODELS", font_size=Pt(26))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Input: 135 features per subject",
            "Models: XGBoost, LightGBM, CatBoost",
            "Protocol: 4-fold GroupKFold",
            "  (per official EMS Train_Valid.xlsx splits)",
            "",
            "Results (seed 42):",
            "  XGBoost:  AUC=0.8702, ACC=78.1%",
            "  LightGBM: AUC=0.8769, ACC=78.1%",
            "  CatBoost: AUC=0.8981, ACC=81.25%",
            "",
            "SHAP analysis: pupil_mean and fix_dur_mean",
            "are the most important features",
            "",
            "OOF predictions → Tier 5 Meta-Learner",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier3_shap_summary.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(3.0))
add_picture_safe(slide, "results/tier3_tabular_comparison.png",
                 Inches(4.9), Inches(4.3), Inches(4.9), Inches(2.0))

add_footer_line(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12: Giao thức đánh giá
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "EVALUATION PROTOCOL", font_size=Pt(26))
bullets = [
    "4-fold GroupKFold from Train_Valid.xlsx (official EMS split):",
    "  → Ensures no data leakage across folds",
    "  → Each fold: ~120 train / ~40 validation subjects",
    "",
    "Out-Of-Fold (OOF) Scheme:",
    "  → Each subject appears exactly once in validation",
    "  → OOF predictions = pseudo-test scores on full training set",
    "  → OOF AUC = primary evaluation metric (unbiased)",
    "",
    "3-seed stability (seeds: 42, 123, 456):",
    "  → Report mean ± std to assess stability",
    "  → Reduces sensitivity to random initialization",
    "",
    "Pupil leakage correction: Fold-wise renormalization of pupil features",
    "  → Prevents information leakage from future fixations",
]
set_content(slide, bullets, font_size=Pt(16), idx=13)
add_footer_line(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 13: Tier 4A - GNN-CEFAM
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "TIER 4A — GNN-CEFAM", font_size=Pt(26))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Spatiotemporal graph construction:",
            "  Each fixation = 1 node",
            "  Edges: k-NN (k=3) + temporal neighbors",
            "  Node features: 2053-dim",
            "    (2048 ResNet50 visual + 5 fixation)",
            "",
            "ST-GNN: GAT layers (Graph Attention Network)",
            "  → Learns spatial correlations between fixations",
            "",
            "CEFAM Fusion:",
            "  Expert stream: 135-dim handcrafted MLP",
            "  Graph stream: ST-GNN output",
            "  Bidirectional cross-attention:",
            "    Expert queries fixation nodes (Direction 1)",
            "    → Enables interpretation via saliency attention maps",
            "",
            "4-fold, pupil leakage corrected",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(13)
        break

add_picture_safe(slide, "results/tier4_graph_topology.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(2.5))
add_picture_safe(slide, "results/tier4_real_attention.png",
                 Inches(4.9), Inches(3.9), Inches(4.9), Inches(2.5))

add_footer_line(slide)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ============================================================
# SLIDE 14: Tier 4B - BiCA-HS
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "TIER 4B — BiCA-HS (BIDIRECTIONAL CROSS-ATTENTION HYBRID STREAM)", font_size=Pt(20))
bullets = [
    "Stream 1 — Transformer Encoder on raw fixation sequences:",
    "  Input: 5-dim fixation sequence [x, y, t, duration, pupil], max_seq_len=200",
    "  → Learns temporal patterns of the full scanpath",
    "",
    "Stream 2 — Biomarker MLP on 15 stimulus-level features:",
    "  d_bio = 15 (fix_count, fix_dur_*, sacc_*, scanpath_length, hull_area, pupil_*...)",
    "  → Learns aggregated biomarker patterns per stimulus",
    "",
    "Bidirectional Cross-Attention (BiCA):",
    "  Q_seq → K_bio, V_bio  (sequence queries biomarker context)",
    "  Q_bio → K_seq, V_seq  (biomarker queries sequence context)",
    "  → Bidirectional fusion of fine-grained sequence and high-level biomarkers",
    "",
    "Training: 4-fold, 150 epochs max, early stopping patience=15, CosineAnnealingLR",
    "Results: AUC=0.9522, ACC=90.6%, Sensitivity=96.25% (seed 42) — best single model",
]
set_content(slide, bullets, font_size=Pt(16), idx=13)
add_footer_line(slide)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ============================================================
# SLIDE 15: Tier 5 - Meta-Learner
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "TIER 5 — META-LEARNER", font_size=Pt(26))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Input: OOF predictions from:",
            "  CatBoost (best tabular model)",
            "  BiCA-HS (best DL model)",
            "",
            "Two combination strategies:",
            "  1. Weighted Average (Optuna)",
            "     → Optimizes weight α ∈ [0,1]",
            "  2. Logistic Regression (4-fold CV)",
            "     → Learns linear combination",
            "",
            "Selection: Compare OOF AUC across strategies",
            "",
            "⚠️ Important note:",
            "  Tier 5 AUC is IN-SAMPLE (training data only)",
            "  Serves as exploratory indicator only",
            "  DeLong test: Tier5 vs BiCA-HS",
            "  Z≈0.6, p≈0.54 (not statistically significant)",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier5_s42/figures/tier5_roc_comparison.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(5.0))

add_footer_line(slide)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ============================================================
# SLIDE 16: Kết quả chính - Bảng so sánh
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "MAIN RESULTS — MODEL COMPARISON (SEED 42)", font_size=Pt(22))

headers = ["Model", "AUC", "ACC", "F1", "Note"]
rows = [
    ["MSNet (SOTA)", "0.8854", "81.25%", "—", "Reference baseline"],
    ["XGBoost", "0.8702", "78.12%", "0.7853", "Tier 3 tabular"],
    ["LightGBM", "0.8769", "78.12%", "0.7898", "Tier 3 tabular"],
    ["CatBoost", "0.8981", "81.25%", "0.8500", "Tier 3, best tabular"],
    ["GNN-CEFAM", "0.9405", "86.88%", "0.8679", "Tier 4A ✓ surpasses SOTA"],
    ["BiCA-HS", "0.9522", "90.62%", "0.9123", "Tier 4B ✓ best single model"],
    ["Tier 5 (exploratory)", "0.9594*", "86.25%", "0.8590", "*In-sample, exploratory"],
]
add_table_slide(
    slide, headers, rows,
    left=Inches(0.2), top=Inches(1.1),
    width=Inches(9.6), height=Inches(4.5),
    font_size=Pt(14)
)

add_textbox(
    slide,
    "BiCA-HS surpasses SOTA: +6.68% AUC, +9.37% ACC | GNN-CEFAM surpasses SOTA: +5.51% AUC, +5.63% ACC | *Tier 5 is exploratory (in-sample)",
    left=Inches(0.2), top=Inches(5.7),
    width=Inches(9.6), height=Inches(0.6),
    font_size=Pt(12), bold=True, color=RED,
    align=PP_ALIGN.CENTER
)

add_footer_line(slide)
add_slide_number(slide, 16, TOTAL_SLIDES)

# ============================================================
# SLIDE 17: ROC Curves + DCA
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "ROC CURVES & DECISION CURVE ANALYSIS (DCA)", font_size=Pt(24))

add_picture_safe(slide, "results/tier5_s42/figures/tier5_roc_comparison.png",
                 Inches(0.2), Inches(1.2), Inches(4.6), Inches(4.7))
add_picture_safe(slide, "results/tier5_s42/figures/tier5_dca.png",
                 Inches(5.0), Inches(1.2), Inches(4.7), Inches(4.7))

add_textbox(
    slide, "ROC comparing 5 models + SOTA | DCA: Net Benefit across threshold probabilities",
    left=Inches(0.2), top=Inches(6.1),
    width=Inches(9.6), height=Inches(0.3),
    font_size=Pt(11), color=DARK_GRAY,
    align=PP_ALIGN.CENTER
)
add_footer_line(slide)
add_slide_number(slide, 17, TOTAL_SLIDES)

# ============================================================
# SLIDE 18: Kiểm định thống kê
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "STATISTICAL TESTING", font_size=Pt(26))

# Bootstrap CI table
add_textbox(slide, "Bootstrap CI (n=10,000) vs. MSNet SOTA (AUC=0.8854):",
            left=Inches(0.3), top=Inches(1.1),
            width=Inches(9.4), height=Inches(0.4),
            font_size=Pt(16), bold=True, color=RED)

headers_bt = ["Model", "Seed 42", "Seed 123", "Seed 456", "Conclusion"]
rows_bt = [
    ["BiCA-HS", "p=0.0002 ✓", "p=0.0162 ✓", "p=0.0011 ✓", "ALL statistically significant"],
    ["GNN-CEFAM", "p=0.0039 ✓", "p=0.0010 ✓", "p=0.0120 ✓", "ALL statistically significant"],
    ["CatBoost", "p≈0.29 ✗", "p≈0.29 ✗", "p≈0.29 ✗", "NOT statistically significant"],
]
add_table_slide(
    slide, headers_bt, rows_bt,
    left=Inches(0.2), top=Inches(1.55),
    width=Inches(9.6), height=Inches(1.8),
    font_size=Pt(13)
)

# DeLong test result
add_textbox(slide, "DeLong Test (Tier 5 vs BiCA-HS):",
            left=Inches(0.3), top=Inches(3.5),
            width=Inches(9.4), height=Inches(0.4),
            font_size=Pt(16), bold=True, color=BLUE_ACCENT)

add_textbox(
    slide,
    "Z ≈ 0.6,  p ≈ 0.54  →  Tier 5 does NOT significantly improve over BiCA-HS alone\n"
    "→ Conclusion: BiCA-HS is the best single model and sufficiently strong",
    left=Inches(0.3), top=Inches(3.95),
    width=Inches(9.4), height=Inches(0.8),
    font_size=Pt(15), color=DARK_GRAY
)

# Multi-seed summary
add_textbox(slide, "Multi-seed (mean ± std, seeds: 42, 123, 456):",
            left=Inches(0.3), top=Inches(4.85),
            width=Inches(9.4), height=Inches(0.4),
            font_size=Pt(14), bold=True, color=DARK_GRAY)
add_textbox(
    slide,
    "BiCA-HS: AUC=0.9454±0.0099, ACC=0.8792±0.0262, Sensitivity=0.9625±0.0177\n"
    "GNN-CEFAM: AUC=0.9398±0.0047, ACC=0.8667±0.0179\n"
    "CatBoost: AUC=0.8990±0.0012 (most stable but lowest)",
    left=Inches(0.3), top=Inches(5.3),
    width=Inches(9.4), height=Inches(0.9),
    font_size=Pt(13), color=DARK_GRAY
)

add_footer_line(slide)
add_slide_number(slide, 18, TOTAL_SLIDES)

# ============================================================
# SLIDE 19: Ablation Study
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "ABLATION STUDY — COMPONENT CONTRIBUTIONS", font_size=Pt(22))

# Left: table
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "Ablation results (seed 42):",
            "",
            "XGBoost (tabular only)",
            "  AUC = 0.8702 — baseline",
            "",
            "ST-GNN (graph only, no CEFAM)",
            "  AUC = 0.9309",
            "  +Graph topology: +0.0607 AUC",
            "",
            "GNN-CEFAM (full Tier 4A)",
            "  AUC = 0.9405",
            "  +CEFAM fusion: +0.0096 AUC",
            "",
            "BiCA-HS (Tier 4B)",
            "  AUC = 0.9522 — best single model",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "experiments/ablation/figures/F2_component_contribution.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(5.0))

add_footer_line(slide)
add_slide_number(slide, 19, TOTAL_SLIDES)

# ============================================================
# SLIDE 20: Phân tích theo danh mục stimulus
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "STIMULUS CATEGORY ANALYSIS", font_size=Pt(24))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "BiCA-HS per-category AUC (seed 42):",
            "",
            "  Social (portraits):      AUC = 0.9523",
            "  Manipulated (edited):    AUC = 0.9523",
            "  Natural (landscapes):    AUC = 0.9513",
            "  Synthetic (textures):    AUC = 0.9450",
            "",
            "Observations:",
            "  Consistent performance across all 4 categories",
            "  Social/Manipulated: highest AUC",
            "    → SZ shows clearer anomalies when viewing people",
            "  Synthetic: lowest but still very high",
            "    → Model generalizes well, not dependent",
            "    on specific semantic features",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "experiments/ablation/figures/F3_category_analysis.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(5.0))

add_footer_line(slide)
add_slide_number(slide, 20, TOTAL_SLIDES)

# ============================================================
# SLIDE 21: Khả năng giải thích - SHAP
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "EXPLAINABILITY — SHAP ANALYSIS", font_size=Pt(24))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "SHAP (SHapley Additive exPlanations):",
            "  Explains each feature's contribution",
            "  for each individual prediction",
            "",
            "Top features by importance (CatBoost):",
            "  1. pupil_mean, pupil_std — pupil size",
            "  2. fix_dur_mean — fixation duration",
            "  3. sacc_amp_mean — saccade amplitude",
            "  4. center_bias — central viewing tendency",
            "  5. spatial_entropy — spatial dispersion",
            "",
            "Clinical findings:",
            "  SZ shows abnormal pupil dilation (larger pupil)",
            "  SZ has shorter fixation duration than HC",
            "  SZ has less structured scanpaths than HC",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier3_shap_summary.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(5.0))

add_footer_line(slide)
add_slide_number(slide, 21, TOTAL_SLIDES)

# ============================================================
# SLIDE 22: Khả năng giải thích - Attention
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "EXPLAINABILITY — ATTENTION VISUALIZATION", font_size=Pt(22))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "CEFAM Attention (GNN-CEFAM):",
            "  Expert queries fixation nodes",
            "  → Identifies the most task-relevant fixations",
            "  → Interpretable by image region",
            "",
            "BiCA-HS Attention:",
            "  Transformer self-attention",
            "  → Represents temporal correlations",
            "  between fixations in the scanpath",
            "",
            "Visualization results:",
            "  SZ: attention concentrates on peripheral regions",
            "  HC: attention follows semantically meaningful regions",
            "  Consistent with clinical theory",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier4_scanpath_attention.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(5.0))

add_footer_line(slide)
add_slide_number(slide, 22, TOTAL_SLIDES)

# ============================================================
# SLIDE 23: Multi-seed Stability
# ============================================================
slide = add_slide(prs, 4)

set_title(slide, "MULTI-SEED STABILITY", font_size=Pt(22))

for ph in slide.placeholders:
    if ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        bullets = [
            "3 seeds: 42, 123, 456",
            "",
            "BiCA-HS per-fold (seed 42):",
            "  Fold 0: AUC = 0.9141",
            "  Fold 1: AUC = 0.9583",
            "  Fold 2: AUC = 1.0000",
            "  Fold 3: AUC = 0.9242",
            "",
            "Fold 2 AUC = 1.0 phenomenon:",
            "  Perfect separation for 40 subjects in Fold 2",
            "  Not global overfitting",
            "  (all other folds are <1.0)",
            "",
            "Multi-seed mean ± std:",
            "  BiCA-HS: AUC=0.9454±0.0099",
            "  GNN-CEFAM: AUC=0.9398±0.0047",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.strip()
            p.level = 1 if b.startswith('  ') else 0
            for run in p.runs:
                run.font.size = Pt(14)
        break

add_picture_safe(slide, "results/tier5_multiseed_stability.png",
                 Inches(4.9), Inches(1.2), Inches(4.9), Inches(2.5))
add_picture_safe(slide, "experiments/ablation/figures/F5_fold_stability.png",
                 Inches(4.9), Inches(3.9), Inches(4.9), Inches(2.5))

add_footer_line(slide)
add_slide_number(slide, 23, TOTAL_SLIDES)

# ============================================================
# SLIDE 24: Hạn chế & Reviewer Concerns
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "LIMITATIONS & POTENTIAL CONCERNS", font_size=Pt(26))
bullets = [
    "1. Delta feature contamination risk:",
    "  → Delta features computed from training fold; fold-wise normalization applied but warrants close verification",
    "",
    "2. BiCA-HS d_bio=15 vs config files:",
    "  → Some config files may record different d_bio values; confirmed at runtime to use 15 stimulus-level features",
    "",
    "3. Fold 2 = 1.0 (perfect AUC) phenomenon:",
    "  → 40 subjects in Fold 2 may be linearly separable; not an error but should be noted in reporting",
    "  → Remaining three folds all <1.0; mean still 0.9454",
    "",
    "4. Tier 5 AUC is in-sample:",
    "  → Cannot be used as primary result; serves as exploratory metric only",
    "  → DeLong test confirms no significant improvement over BiCA-HS",
    "",
    "5. Test set (48 subjects): labels withheld — final test results not yet available",
]
set_content(slide, bullets, font_size=Pt(15), idx=13)
add_footer_line(slide)
add_slide_number(slide, 24, TOTAL_SLIDES)

# ============================================================
# SLIDE 25: Kết luận & Hướng phát triển
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "CONCLUSION & FUTURE WORK", font_size=Pt(26))
bullets = [
    "Conclusions — Answering research questions:",
    "  RQ1: 5-tier pipeline effectively integrates features across multiple granularity levels",
    "  RQ2: CEFAM and BiCA-HS surpass SOTA with statistical significance (bootstrap CI, p<0.05)",
    "  RQ3: Delta features improve SZ/HC discrimination across stimulus categories",
    "  RQ4: BiCA-HS is the best single model: AUC=0.9454±0.0099",
    "",
    "Contribution summary:",
    "  ✓ Most comprehensive pipeline on the EMS dataset",
    "  ✓ BiCA-HS: +6.68% AUC over SOTA MSNet",
    "  ✓ Interpretable via SHAP + attention visualization",
    "",
    "Future directions:",
    "  → Transfer learning from other eye-tracking datasets",
    "  → Self-supervised pre-training on fixation sequences",
    "  → Integrate additional modalities: EEG, fMRI with eye-tracking",
    "  → Real-world deployment: real-time screening tool",
    "  → Evaluation on test set once labels are released",
]
set_content(slide, bullets, font_size=Pt(15), idx=13)
add_footer_line(slide)
add_slide_number(slide, 25, TOTAL_SLIDES)

# ============================================================
# SLIDE 26: Tài liệu tham khảo
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "REFERENCES", font_size=Pt(26))
refs = [
    "[1] Song et al. (2025). EMS: Eye Movement-based Schizophrenia Recognition Dataset. arXiv:2502.xxxxx",
    "[2] MSNet: Multi-Scale Network for schizophrenia detection from eye movements. AUC=0.8854.",
    "[3] Velickovic et al. (2018). Graph Attention Networks. ICLR 2018.",
    "[4] Vaswani et al. (2017). Attention is All You Need. NeurIPS 2017.",
    "[5] Chen & Guestrin (2016). XGBoost: A Scalable Tree Boosting System. KDD 2016.",
    "[6] Prokhorenkova et al. (2018). CatBoost: unbiased boosting with categorical features. NeurIPS 2018.",
    "[7] Lundberg & Lee (2017). A unified approach to interpreting model predictions (SHAP). NeurIPS 2017.",
    "[8] He et al. (2016). Deep Residual Learning for Image Recognition (ResNet). CVPR 2016.",
    "[9] Vickers & Elkin (2006). Decision curve analysis: a novel method for evaluating prediction models. MDM.",
    "[10] DeLong et al. (1988). Comparing the areas under two or more correlated receiver operating curves. Biometrics.",
]
set_content(slide, refs, font_size=Pt(14), idx=13)
add_footer_line(slide)
add_slide_number(slide, 26, TOTAL_SLIDES)

# ============================================================
# SLIDE 27: Q&A
# ============================================================
slide = add_slide(prs, 3)

set_title(slide, "", font_size=Pt(28))

add_rect(slide, 0, 0, SW, SH, fill_color=RED)

add_textbox(
    slide,
    "THANK YOU FOR YOUR ATTENTION!",
    left=Inches(0.5), top=Inches(2.0),
    width=SW - Inches(1.0), height=Inches(1.0),
    font_size=Pt(32), bold=True, color=WHITE,
    align=PP_ALIGN.CENTER
)
add_textbox(
    slide,
    "QUESTIONS & DISCUSSION",
    left=Inches(0.5), top=Inches(3.1),
    width=SW - Inches(1.0), height=Inches(0.8),
    font_size=Pt(28), bold=False, color=WHITE,
    align=PP_ALIGN.CENTER
)
add_textbox(
    slide,
    "Thesis: Eye Movement-Based Schizophrenia Recognition\nusing Hierarchical Multi-Stream Learning\nDefense Academy of Vietnam (DAV) — 2025",
    left=Inches(0.5), top=Inches(4.3),
    width=SW - Inches(1.0), height=Inches(1.2),
    font_size=Pt(16), bold=False, color=WHITE,
    align=PP_ALIGN.CENTER
)

add_slide_number(slide, 27, TOTAL_SLIDES)

# ============================================================
# SAVE
# ============================================================
print("\nSaving presentation to:", OUTPUT)
prs.save(OUTPUT)
print("Saved successfully!")

# Verify file size
size = os.path.getsize(OUTPUT)
print("File size: %d bytes (%.1f MB)" % (size, size / 1024 / 1024))
if size > 500000:
    print("OK: File size is > 500KB as expected.")
else:
    print("WARNING: File size is smaller than expected!")

print("\nTotal slides created:", len(prs.slides))
print("DONE.")
