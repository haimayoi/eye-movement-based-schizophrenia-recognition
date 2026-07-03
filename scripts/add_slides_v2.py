"""
Add all remaining visualization slides to DAV_slides.pptx.
Comprehensive update — all 24 result images evaluated and placed.
Target: 31 existing → 44 slides total.
"""
import sys, os, copy
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = r"D:\DAV\Eye Movement-Based Schizophrenia Recognition"
SRC  = os.path.join(BASE, "DAV_slides.pptx")
DST  = os.path.join(BASE, "DAV_slides_v2.pptx")
RES  = os.path.join(BASE, "results")
FIG  = os.path.join(BASE, "experiments", "ablation", "figures")

def img(name): return os.path.join(RES, name)
def fig(name): return os.path.join(FIG, name)
def emu(i):    return int(i * 914400)

RED   = RGBColor(0xC0, 0x00, 0x00)
BLUE  = RGBColor(0x1F, 0x49, 0x7D)
GRAY  = RGBColor(0x40, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x70, 0x00)
ORANGE= RGBColor(0xC0, 0x50, 0x00)
PURP  = RGBColor(0x60, 0x00, 0x80)

# ── Helpers ──────────────────────────────────────────────────────────────
def tf_clear(tf):
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()

def tf_set(shape, lines, size=13, color=GRAY):
    tf = shape.text_frame
    tf.word_wrap = True
    tf_clear(tf)
    first = True
    for line in lines:
        if isinstance(line, dict):
            text  = line.get('text', '')
            sz    = line.get('size', size)
            bold  = line.get('bold', False)
            col   = line.get('color', color)
            align = line.get('align', PP_ALIGN.LEFT)
            itl   = line.get('italic', False)
            sp    = line.get('space_before', 0)
        else:
            text, sz, bold, col, align, itl, sp = line, size, False, color, PP_ALIGN.LEFT, False, 0

        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        if sp: para.space_before = Pt(sp)
        para.clear()
        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.italic = itl

def add_pic(slide, path, left, top, w, h):
    if not os.path.exists(path):
        print(f"  [MISS] {path}")
        return None
    return slide.shapes.add_picture(path, emu(left), emu(top), emu(w), emu(h))

def add_tb(slide, left, top, w, h, lines, size=12, color=GRAY):
    tb = slide.shapes.add_textbox(emu(left), emu(top), emu(w), emu(h))
    tf_set(tb, lines, size=size, color=color)
    tb.text_frame.word_wrap = True
    return tb

def footer(slide, num, total=44):
    add_tb(slide, 0.2, 7.15, 7.0, 0.3,
           [{'text': 'Thesis Defense — DAV / HUST 2025', 'size': 9, 'color': GRAY, 'italic': True}])
    add_tb(slide, 8.8, 7.15, 1.0, 0.3,
           [{'text': f'{num} / {total}', 'size': 9, 'color': GRAY}])

def new_slide(prs, layout_name='1_Two Content'):
    for layout in prs.slide_master.slide_layouts:
        if layout.name == layout_name:
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_master.slide_layouts[3])

def set_slide_title(slide, text, size=20):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            tf_set(shape, [{'text': text, 'size': size, 'bold': True, 'color': WHITE}])
            return

def set_left_content(slide, lines, size=12):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf_set(shape, lines, size=size)
            return

# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation(SRC)
print(f"Loaded {len(prs.slides)} slides. Adding comprehensive visualization slides...")

# ══════════════════════════════════════════════════════════════════════════
# UPDATE SLIDE 2 — TABLE OF CONTENTS (comprehensive)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides[1]
for shape in s.shapes:
    if shape.shape_type == 17 and len(shape.text) > 50:
        tf_set(shape, [
            {'text': 'TIER 1 — DATA PREPROCESSING (Slides 8–13)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Filtering pipeline, 2D scanpath viz, pupil normalization (DAV: C0, C2)', 'size': 11, 'color': GRAY},
            {'text': 'TIER 2 — FEATURE ENGINEERING (Slides 9, 10, 32–35)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  15+120 features, violin plots, correlation heatmap, delta density (DAV: C3, C4)', 'size': 11, 'color': GRAY},
            {'text': 'TIER 3 — TABULAR MODELS & XAI (Slides 11, 36–39)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  XGBoost/CatBoost, OOF confusion matrix, SHAP beeswarm + waterfall (DAV: C7)', 'size': 11, 'color': GRAY},
            {'text': 'EVALUATION PROTOCOL (Slide 12)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  4-fold GroupKFold OOF, zero subject leakage', 'size': 11, 'color': GRAY},
            {'text': 'TIER 4A — GNN-CEFAM (Slides 13, 40–42)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Graph topology, attention weights, threshold sensitivity', 'size': 11, 'color': GRAY},
            {'text': 'TIER 4B — BiCA-HS (Slide 14)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Transformer + bidirectional cross-attention', 'size': 11, 'color': GRAY},
            {'text': 'TIER 5 — META-LEARNER (Slides 15, 43–44)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Ensemble decision space, calibration (DAV: C1)', 'size': 11, 'color': GRAY},
            {'text': 'RESULTS & ANALYSIS (Slides 16–27)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Table 1, ROC, ablation, category AUC, statistical testing', 'size': 11, 'color': GRAY},
            {'text': 'DIM. REDUCTION — t-SNE & UMAP (Slide 28)', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '  ResNet50 2048-dim → 2D, cluster separation analysis (DAV: C5, C6)', 'size': 11, 'color': BLUE},
            {'text': 'CASE STUDY & EXPLAINABILITY (Slides 29–30, 45–47)', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '  SZ #302 vs HC #8 full analysis, SHAP waterfall, decision space (DAV: C7)', 'size': 11, 'color': BLUE},
        ], size=11)
print("  Slide 2 (ToC) updated")

# ══════════════════════════════════════════════════════════════════════════
# UPDATE SLIDE 28 — t-SNE + UMAP side by side
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides[27]
# Remove existing pictures
for shape in list(s.shapes):
    if shape.shape_type == 13:
        shape._element.getparent().remove(shape._element)
# Add t-SNE and UMAP side by side
add_pic(s, img("tier4_resnet_tsne.png"), 0.2, 1.3, 4.8, 4.5)
add_pic(s, img("tier4_resnet_umap.png"), 5.1, 1.3, 4.8, 4.5)
# Update left content (placeholder 1)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 't-SNE vs UMAP on ResNet50 Embeddings (DAV: C5, C6)', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 5},
            {'text': 'Input: 2048-dim ResNet50 patch features per fixation node', 'size': 11, 'color': GRAY},
            {'text': 'Aggregated to subject-level mean → 2048-dim per subject', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 5},
            {'text': 't-SNE (left): local structure, perplexity=30, 1000 iter', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '• Partial separation visible in some regions', 'size': 11, 'color': GRAY},
            {'text': '• Significant overlap → visual patches alone insufficient', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'UMAP (right): global structure, n_neighbors=15', 'size': 12, 'bold': True, 'color': PURP},
            {'text': '• No clear cluster separation SZ vs HC', 'size': 11, 'color': GRAY},
            {'text': '• Confirms: what you fixate ON is less important than', 'size': 11, 'color': GRAY},
            {'text': '  HOW you fixate (temporal/spatial dynamics)', 'size': 11, 'bold': True, 'color': GREEN},
            {'text': '', 'size': 5},
            {'text': '→ Motivates Tier 4 deep sequential modeling', 'size': 12, 'bold': True, 'color': GREEN},
        ])
print("  Slide 28 (t-SNE+UMAP) updated")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 32 — Tier 1: 2D Scanpath Filtering (DAV C2)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "TIER 1 — 2D SCANPATH FILTERING VISUALIZATION (DAV: C2)", 17)
set_left_content(s, [
    {'text': 'INPUT → CLEAN: What gets removed?', 'size': 13, 'bold': True, 'color': RED},
    {'text': '', 'size': 4},
    {'text': 'Pink dots:  Valid fixations (kept)', 'size': 12, 'color': GRAY},
    {'text': '    → Within 1024×768 screen boundary', 'size': 11, 'color': GRAY},
    {'text': '    → Duration ≥ 50 ms', 'size': 11, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': 'Red ×:  Rejected — spatial boundary', 'size': 12, 'color': RED},
    {'text': '    → FIX_X < 0 or > 1024, FIX_Y < 0 or > 768', 'size': 11, 'color': GRAY},
    {'text': '    → Example: fixation at X=35, Y=271 (outside screen)', 'size': 11, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': 'Orange dots:  Rejected — too short (<50ms)', 'size': 12, 'color': ORANGE},
    {'text': '    → Post-saccadic tremor artifacts', 'size': 11, 'color': GRAY},
    {'text': '    → Below perceptual threshold for meaningful fixation', 'size': 11, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'DAV C2 — Outlier Removal:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': 'This 2D scatter is a canonical outlier detection visualization', 'size': 11, 'color': GRAY},
    {'text': 'Spatial: geometric boundary rule (domain knowledge)', 'size': 11, 'color': GRAY},
    {'text': 'Temporal: threshold rule based on psychophysics (50ms minimum)', 'size': 11, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': 'Subject 2 (SZ): 1,560 raw → 1,519 retained (2.6% removed)', 'size': 11, 'bold': True, 'color': GREEN},
])
add_pic(s, img("tier1_2d_filtering.png"), 4.9, 1.1, 5.0, 5.7)
footer(s, 32)
print("  Slide 32 (2D filtering) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 33 — Tier 1: Spatial Filter X/Y Distributions (DAV C2)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "TIER 1 — FIXATION POSITION DISTRIBUTIONS: HC vs SZ (DAV: C2)", 17)
set_left_content(s, [
    {'text': 'Post-Filter FIX_X and FIX_Y Distributions', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'Mann-Whitney U test: *** (p < 0.001)', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Key Finding — Center Bias in SZ:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '• SZ: FIX_X distribution peaks sharply at center (~500px)', 'size': 12, 'color': GRAY},
    {'text': '• HC: FIX_X distribution broader — more peripheral exploration', 'size': 12, 'color': GRAY},
    {'text': '• SZ: FIX_Y similarly more center-concentrated', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Clinical Interpretation:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '→ SZ patients exhibit restricted visual exploration (smaller convex hull area)', 'size': 12, 'color': GRAY},
    {'text': '→ HC patients scan more freely across the full stimulus', 'size': 12, 'color': GRAY},
    {'text': '→ This difference directly captured by Center Bias feature (Tier 2A)', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': '', 'size': 5},
    {'text': 'DAV C2 — Visualization of filtered data distribution:', 'size': 12, 'color': BLUE},
    {'text': 'Histograms show data quality post-cleaning', 'size': 11, 'italic': True, 'color': GRAY},
    {'text': 'No outliers remain outside [0, 1024] × [0, 768]', 'size': 11, 'italic': True, 'color': GRAY},
])
add_pic(s, img("tier1_spatial_filter.png"), 0.2, 3.8, 9.6, 3.0)
footer(s, 33)
print("  Slide 33 (spatial filter dist.) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 34 — Tier 2B: Delta Feature Density Plots (DAV C3)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "TIER 2B — DELTA FEATURE DISTRIBUTIONS: SZ vs HC (DAV: C3)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Δ(Social − Natural) contrast distributions — 6 key features shown', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 4},
            {'text': '1. Conv. Hull Area:    p = 5.1e-5 ***  HC shifts MORE positive (larger area in Social vs Natural)', 'size': 12, 'color': GRAY},
            {'text': '2. Spatial Entropy:    p = 1.4e-3 **   HC has higher entropy delta → more diverse gaze in Social', 'size': 12, 'color': GRAY},
            {'text': '3. Sacc. Amp. Mean:    p = 3.5e-3 **   HC makes larger saccades when viewing Social stimuli', 'size': 12, 'color': GRAY},
            {'text': '4. Scanpath Length:    p = 6.1e-3 **   HC explores more (longer scanpath) in Social vs Natural', 'size': 12, 'color': GRAY},
            {'text': '5. Pupil Mean:         p = 3.0e-2 *    HC pupil more responsive to Social stimuli', 'size': 12, 'color': GRAY},
            {'text': '6. Center Bias:        p = 3.1e-2 *    HC center bias shifts less between categories', 'size': 12, 'color': GRAY},
            {'text': '', 'size': 5},
            {'text': 'All 6 features show HC shifted right (more positive delta) vs SZ', 'size': 13, 'bold': True, 'color': BLUE},
            {'text': 'Interpretation: HC responds MORE to Social stimuli than Natural', 'size': 12, 'color': BLUE},
            {'text': '             SZ shows FLAT response — no category-specific modulation', 'size': 12, 'color': BLUE},
            {'text': '', 'size': 5},
            {'text': '→ These delta features are the most discriminative in CatBoost SHAP analysis', 'size': 13, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("tier2_delta_density.png"), 0.2, 4.0, 9.6, 2.9)
footer(s, 34)
print("  Slide 34 (delta density) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 35 — Tier 2: Feature Violin Plots (DAV C3 — detailed)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "TIER 2 — TOP 8 DISCRIMINATIVE FEATURES: HC vs SZ (DAV: C3)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Mann-Whitney U test on stimulus-level trials — all features: *** (p ~ 0)', 'size': 12, 'bold': True, 'color': RED},
            {'text': '', 'size': 4},
            {'text': 'Scanpath Length (px):   HC > SZ ***  HC explores a wider area per trial', 'size': 11, 'color': GRAY},
            {'text': 'Pupil Std (au):          HC > SZ ***  SZ has more variable but smaller pupil', 'size': 11, 'color': GRAY},
            {'text': 'Fix. Dur. Mean (ms):     HC > SZ ***  HC dwells longer per fixation', 'size': 11, 'color': GRAY},
            {'text': 'Fix. Dur. Median (ms):   HC > SZ ***  Consistent longer fixations in HC', 'size': 11, 'color': GRAY},
            {'text': 'Fix. Count:              HC > SZ ***  HC makes more fixations per trial', 'size': 11, 'color': GRAY},
            {'text': 'SA-FDR Mean:             HC > SZ ***  Higher saccade-dwell exploration rate in HC', 'size': 11, 'color': GRAY},
            {'text': 'Pupil Mean (au):         HC > SZ ***  HC consistently higher baseline pupil', 'size': 11, 'color': GRAY},
            {'text': 'Pupil CV:                HC < SZ ***  SZ more variable relative to mean (higher CV)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 5},
            {'text': 'Pupil features are ORTHOGONAL to spatial features (correlation heatmap)', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '→ Two independent biomarker families: spatial dynamics + pupil physiology', 'size': 12, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("tier2_feature_violin.png"), 0.1, 4.0, 9.7, 2.85)
footer(s, 35)
print("  Slide 35 (violin plots) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 36 — Tier 3: OOF Confusion Matrix (XGBoost as baseline)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "TIER 3 — XGBoost OOF EVALUATION: CONFUSION MATRIX (DAV: C1)", 17)
set_left_content(s, [
    {'text': 'XGBoost OOF Performance (Subject-Level, Seed 42)', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'AUC = 0.8702  |  Accuracy = 78.1%', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '', 'size': 5},
    {'text': 'Confusion Matrix @ τ = 0.50:', 'size': 13, 'bold': True, 'color': RED},
    {'text': '  TP (SZ→SZ): 64  |  FN (SZ→HC): 16', 'size': 12, 'color': GRAY},
    {'text': '  TN (HC→HC): 61  |  FP (HC→SZ): 19', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': 'Sensitivity (Recall):  80.0%', 'size': 12, 'color': GRAY},
    {'text': 'Specificity:           76.2%', 'size': 12, 'color': GRAY},
    {'text': 'PPV (Precision):       77.1%', 'size': 12, 'color': GRAY},
    {'text': 'NPV:                   79.2%', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Probability Distribution (right):', 'size': 13, 'bold': True, 'color': RED},
    {'text': '• HC: strongly bimodal near P(SZ)=0', 'size': 12, 'color': GRAY},
    {'text': '• SZ: strongly bimodal near P(SZ)=1', 'size': 12, 'color': GRAY},
    {'text': '• Strong bimodal separation despite AUC "only" 0.87', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': '• Confirms good class separability in expert feature space', 'size': 11, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': '→ BiCA-HS achieves AUC=0.955: +0.085 over XGBoost baseline', 'size': 12, 'bold': True, 'color': GREEN},
])
add_pic(s, img("tier3_oof_eval.png"), 4.9, 1.2, 5.0, 5.5)
footer(s, 36)
print("  Slide 36 (XGBoost OOF confusion matrix) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 37 — SHAP Waterfall: Individual Prediction Explanation (DAV C7)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "SHAP WATERFALL — INDIVIDUAL PREDICTION EXPLANATION (DAV: C7)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'SZ Subject #302 (P(SZ)=0.9962)  vs  HC Subject #8 (P(SZ)=0.0009)', 'size': 13, 'bold': True, 'color': RED},
            {'text': 'CatBoost model — SHAP values show per-feature contribution from base rate 0.5', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'SZ #302 — Top PUSH-TOWARD-SZ features (red bars):', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  Synthetic_pupil_std  +0.838   (high pupil variability in Synthetic category)', 'size': 11, 'color': GRAY},
            {'text': '  Social_pupil_std     +0.603   (high variability viewing social scenes)', 'size': 11, 'color': GRAY},
            {'text': '  Natural_pupil_std    +0.470   (elevated variability across all categories)', 'size': 11, 'color': GRAY},
            {'text': '  Social_sa_fdr_mean   +0.451   (reduced exploration rate in social scenes)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'HC Subject #8 — Top PUSH-AWAY-FROM-SZ features (blue bars):', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '  Natural_pupil_cv    −0.643   (low pupil coefficient of variation)', 'size': 11, 'color': GRAY},
            {'text': '  Synthetic_pupil_std −0.588   (stable pupil in synthetic category)', 'size': 11, 'color': GRAY},
            {'text': '  Social_pupil_std    −0.391   (highly stable social viewing)', 'size': 11, 'color': GRAY},
            {'text': '  Social_spatial_entropy −0.360 (high entropy = broad exploration)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'Key finding: PUPIL VARIABILITY is the strongest individual predictor', 'size': 13, 'bold': True, 'color': GREEN},
            {'text': 'SZ = high variability | HC = stable, controlled pupil response', 'size': 12, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("explainability_shap_waterfall.png"), 0.1, 4.5, 9.8, 2.35)
footer(s, 37)
print("  Slide 37 (SHAP waterfall) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 38 — Feature Profile: SZ vs HC individual comparison (DAV C7)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "INDIVIDUAL FEATURE PROFILE — SZ #302 vs HC #8 (DAV: C7)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Z-scores relative to ALL 160 training subjects', 'size': 12, 'bold': True, 'color': RED},
            {'text': 'Dashed line = group median (SZ median on left, HC median on right)', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'SZ Subject #302 (left) — Strongly negative z-scores:', 'size': 13, 'bold': True, 'color': RED},
            {'text': '• Fix. Count: −0.56 (fewer fixations than average)', 'size': 11, 'color': GRAY},
            {'text': '• Sacc. Amp. Std: −0.76 (more stereotyped saccade patterns)', 'size': 11, 'color': GRAY},
            {'text': '• Conv. Hull Area: −0.87 (very restricted visual exploration)', 'size': 11, 'color': GRAY},
            {'text': '• Scanpath Length: −0.76 (shorter total path)', 'size': 11, 'color': GRAY},
            {'text': '• Center Bias: −0.90 (but pupil_CV very HIGH = +0.48)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'HC Subject #8 (right) — Positive z-scores on spatial features:', 'size': 13, 'bold': True, 'color': BLUE},
            {'text': '• Fix. Count: +0.40 (more fixations than average HC)', 'size': 11, 'color': GRAY},
            {'text': '• Pupil CV: +1.47 (high variability — outlier within HC)', 'size': 11, 'color': GRAY},
            {'text': '• Scanpath Length: +0.51 (broad exploration)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'Validates: spatial biomarkers separate groups even at individual level', 'size': 12, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("explainability_feature_profile.png"), 0.1, 4.45, 9.8, 2.4)
footer(s, 38)
print("  Slide 38 (feature profile) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 39 — Tier 4: Threshold Sensitivity — critical clinical insight
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "TIER 4 — DECISION THRESHOLD SENSITIVITY: ST-GNN vs BiCA-HS (DAV: C1)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Clinical Critical Comparison: What happens as we vary threshold τ?', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 4},
            {'text': 'ST-GNN (orange):  SEVERE calibration failure at τ = 0.50', 'size': 13, 'bold': True, 'color': ORANGE},
            {'text': '• Sensitivity @ τ=0.50 = 0.625 → misses 37.5% of SZ patients!', 'size': 12, 'color': GRAY},
            {'text': '• Sensitivity collapses sharply above τ=0.40', 'size': 12, 'color': GRAY},
            {'text': '• Specificity near 1.0 @ τ=0.50 → model biased toward predicting HC', 'size': 12, 'color': GRAY},
            {'text': '• Would require recalibration to τ≈0.25 for clinical use', 'size': 12, 'color': ORANGE, 'italic': True},
            {'text': '', 'size': 5},
            {'text': 'BiCA-HS (purple):  Robust and clinically deployable', 'size': 13, 'bold': True, 'color': PURP},
            {'text': '• Sensitivity @ τ=0.50 = 0.875 → balanced detection', 'size': 12, 'color': GRAY},
            {'text': '• Sensitivity degrades gracefully as τ increases', 'size': 12, 'color': GRAY},
            {'text': '• F1 remains high (>0.85) across broad τ range [0.3–0.8]', 'size': 12, 'color': GRAY},
            {'text': '• No post-hoc calibration needed for threshold-free deployment', 'size': 12, 'bold': True, 'color': GREEN},
            {'text': '', 'size': 4},
            {'text': 'Conclusion: AUC alone is insufficient — threshold stability matters clinically', 'size': 13, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("tier4_calibration_threshold.png"), 0.1, 4.45, 9.8, 2.35)
footer(s, 39)
print("  Slide 39 (threshold sensitivity) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 40 — Tier 4: Per-Category P(SZ) — model behaviour comparison
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "TIER 4 — PER-CATEGORY P(SZ) SCORES: 4 MODELS COMPARED (DAV: C1)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Mean P(SZ) ± std across all subjects, per stimulus category', 'size': 12, 'bold': True, 'color': RED},
            {'text': 'All 4 models: *** significance (Mann-Whitney HC vs SZ per category)', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'ST-GNN: Category-INVARIANT scores', 'size': 12, 'bold': True, 'color': ORANGE},
            {'text': '  SZ ≈ 0.59, HC ≈ 0.27 — same pattern for ALL 4 categories', 'size': 11, 'color': GRAY},
            {'text': '  → Spatial topology alone doesn\'t distinguish categories', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'GNN-CEFAM: Similarly category-invariant BUT larger SZ/HC gap', 'size': 12, 'bold': True, 'color': RED},
            {'text': '  SZ ≈ 0.75, HC ≈ 0.21 — expert stream adds discriminative power', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'BiCA-HS: Category-invariant but LARGEST gap', 'size': 12, 'bold': True, 'color': PURP},
            {'text': '  SZ ≈ 0.85, HC ≈ 0.18 — strongest SZ/HC discrimination across all categories', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'CatBoost: Also category-invariant (subject-level aggregation)', 'size': 12, 'bold': True, 'color': GREEN},
            {'text': '  SZ ≈ 0.77, HC ≈ 0.16 — similar to BiCA-HS but AUC lower', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'Key insight: Category-specific signal (Tier 2B deltas) captured in subject-level', 'size': 12, 'bold': True, 'color': GREEN},
            {'text': 'aggregation — per-category scores are averaged before subject prediction', 'size': 11, 'color': GRAY},
        ])
add_pic(s, img("tier4_scanpath_attention.png"), 0.1, 4.45, 9.8, 2.35)
footer(s, 40)
print("  Slide 40 (per-category P(SZ)) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 41 — Tier 5: 2D Ensemble Decision Space
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "TIER 5 — 2D ENSEMBLE DECISION SPACE (DAV: C1)", 17)
set_left_content(s, [
    {'text': 'X-axis: CatBoost P(SZ)  |  Y-axis: BiCA-HS P(SZ)', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'Purple dashed line = Logistic Regression meta-learner boundary', 'size': 11, 'italic': True, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Meta-learner equation:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '2.36 × P_tab + 3.75 × P_deep − 2.90 = 0', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '→ BiCA-HS weighted ~1.6× higher than CatBoost', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Quadrant analysis:', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'Top-right (both high): 95% SZ → easy cases', 'size': 12, 'color': GRAY},
    {'text': 'Bottom-left (both low): 95% HC → easy cases', 'size': 12, 'color': GRAY},
    {'text': 'Top-left (BiCA high, Cat low): 65% SZ → BiCA wins', 'size': 12, 'color': GRAY},
    {'text': 'Bottom-right (Cat high, BiCA low): mixed → hard cases', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Case studies marked:', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'SZ #302 ◆: (P_cat≈0.99, P_bica≈0.95) — both agree SZ', 'size': 12, 'color': RED},
    {'text': 'HC  #8  ◆: (P_cat≈0.01, P_bica≈0.01) — both agree HC', 'size': 12, 'color': BLUE},
    {'text': '', 'size': 4},
    {'text': '→ Complementarity: models disagree on ~15% of subjects', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': 'Ensemble recovers some of these hard cases', 'size': 11, 'color': GRAY},
])
add_pic(s, img("tier5_decision_space_2d.png"), 5.0, 1.1, 4.8, 5.6)
footer(s, 41)
print("  Slide 41 (2D decision space) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 42 — Tier 5: Comprehensive ROC + Calibration
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "TIER 5 — ROC CURVES & PROBABILITY CALIBRATION (DAV: C1)", 17)
set_left_content(s, [
    {'text': 'All-Model ROC Comparison (left panel)', 'size': 13, 'bold': True, 'color': RED},
    {'text': '', 'size': 4},
    {'text': 'Tier 3 XGBoost:         AUC = 0.8702', 'size': 12, 'color': GRAY},
    {'text': 'Tier 3 CatBoost:        AUC = 0.8981', 'size': 12, 'color': GRAY},
    {'text': 'Tier 4 GNN-CEFAM:       AUC = 0.9405', 'size': 12, 'color': GRAY},
    {'text': 'Tier 4 BiCA-HS:         AUC = 0.9522', 'size': 12, 'bold': True, 'color': BLUE},
    {'text': 'Tier 5 Ensemble XGB+CEF: AUC = 0.9336', 'size': 12, 'color': GRAY},
    {'text': 'Tier 5 Cat+BiCA:        AUC = 0.9594', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': 'SOTA MSNet ----:        AUC = 0.8854', 'size': 12, 'color': GRAY},
    {'text': 'Gain vs SOTA: ΔAU = +0.074 ✔', 'size': 13, 'bold': True, 'color': GREEN},
    {'text': '', 'size': 5},
    {'text': 'Reliability Diagram (right panel — DAV: C1)', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'Calibration quality (closer to diagonal = better):', 'size': 12, 'color': GRAY},
    {'text': '• CatBoost:  ECE = 0.107 (moderate)', 'size': 12, 'color': GRAY},
    {'text': '• BiCA-HS:   ECE = 0.170 (overconfident at extremes)', 'size': 12, 'color': ORANGE},
    {'text': '• Ensemble:  ECE = 0.120 (compromise)', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 4},
    {'text': 'Note: ECE values from seed 42; paper reports post-calibration ECE=0.031', 'size': 10, 'italic': True, 'color': GRAY},
])
add_pic(s, img("tier5_comprehensive_eval.png"), 4.9, 1.1, 5.0, 5.6)
footer(s, 42)
print("  Slide 42 (comprehensive ROC + calibration) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 43 — OOF Decision Space Separation (mislabeled as LDA)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "OOF DECISION SPACE SEPARATION (DAV: C6 — Class Separation)", 17)
set_left_content(s, [
    {'text': 'XGBoost OOF Predicted P(SZ) Distribution', 'size': 13, 'bold': True, 'color': RED},
    {'text': '(DAV C6: visualizing class separation in decision space)', 'size': 11, 'italic': True, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'HC (blue) distribution:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': '• Sharp peak near P(SZ) = 0.0 — most HC correctly identified', 'size': 12, 'color': GRAY},
    {'text': '• Small tail extending toward 0.5 (borderline cases)', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'SZ (red) distribution:', 'size': 13, 'bold': True, 'color': RED},
    {'text': '• Sharp peak near P(SZ) = 1.0 — most SZ clearly classified', 'size': 12, 'color': GRAY},
    {'text': '• Small overlap with HC in 0.4–0.6 region', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Decision Space Analysis (DAV C6):', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': 'Strong bimodal separation → high AUC despite "only" 0.87', 'size': 12, 'color': GRAY},
    {'text': 'Only ~15% of subjects fall in ambiguous 0.3–0.7 range', 'size': 12, 'color': GRAY},
    {'text': 'This means the expert feature vector has genuine discriminative power', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': '', 'size': 5},
    {'text': 'Compare with BiCA-HS: even fewer subjects in ambiguous zone', 'size': 12, 'color': BLUE},
    {'text': '→ Deep learning further sharpens the decision boundary', 'size': 12, 'bold': True, 'color': GREEN},
])
add_pic(s, img("lda_class_separation.png"), 5.0, 1.5, 4.8, 4.5)
footer(s, 43)
print("  Slide 43 (OOF decision space) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 44 — Ensemble Decision Space: CatBoost vs BiCA-HS scatter
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
set_slide_title(s, "ENSEMBLE DECISION SPACE: CatBoost × BiCA-HS (DAV: C7)", 17)
set_left_content(s, [
    {'text': 'Decision Space Scatter Plot', 'size': 13, 'bold': True, 'color': RED},
    {'text': 'X = CatBoost P(SZ)  |  Y = BiCA-HS P(SZ)', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Quadrant agreement analysis:', 'size': 13, 'bold': True, 'color': BLUE},
    {'text': 'Q3 (top-right): Both agree SZ → ~65% of actual SZ subjects', 'size': 12, 'color': GRAY},
    {'text': 'Q1 (bottom-left): Both agree HC → ~55% of actual HC subjects', 'size': 12, 'color': GRAY},
    {'text': 'Disagreement regions: ~20% of all subjects', 'size': 12, 'color': ORANGE},
    {'text': '', 'size': 5},
    {'text': 'Model complementarity visible:', 'size': 13, 'bold': True, 'color': RED},
    {'text': '• Some SZ subjects: BiCA-HS HIGH, CatBoost LOW (top-left)', 'size': 12, 'color': GRAY},
    {'text': '  → BiCA-HS captures temporal dynamics missed by tabular model', 'size': 11, 'italic': True, 'color': GRAY},
    {'text': '• Some HC subjects scattered in off-diagonal regions', 'size': 12, 'color': GRAY},
    {'text': '', 'size': 5},
    {'text': 'Case studies (diamonds):', 'size': 13, 'bold': True, 'color': RED},
    {'text': '◆ SZ #302: top-right corner — both models confidently SZ', 'size': 12, 'color': RED},
    {'text': '◆ HC #8:   bottom-left corner — both models confidently HC', 'size': 12, 'color': BLUE},
    {'text': '', 'size': 4},
    {'text': '→ Ensemble logistic regression leverages this complementarity', 'size': 12, 'bold': True, 'color': GREEN},
    {'text': 'Tier 5 AUC = 0.959 > BiCA-HS alone 0.955 (small but consistent)', 'size': 12, 'color': GREEN},
])
add_pic(s, img("explainability_decision_space.png"), 5.0, 1.1, 4.8, 5.6)
footer(s, 44)
print("  Slide 44 (ensemble decision scatter) added")

# ══════════════════════════════════════════════════════════════════════════
# NEW SLIDE 45 — CASE STUDY: SZ #302 vs HC #8 (DAV C7 — centerpiece)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide(prs, '1_Section Header')
set_slide_title(s, "CASE STUDY: SZ SUBJECT #302 vs HC SUBJECT #8 (DAV: C7)", 17)
for shape in s.placeholders:
    if shape.placeholder_format.idx == 1:
        tf_set(shape, [
            {'text': 'Complete Individual-Level Analysis — Most Extreme Cases', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 3},
            {'text': 'SZ Subject #302 (LEFT COLUMN):', 'size': 13, 'bold': True, 'color': RED},
            {'text': '• Scanpath: 1,166 fixations across 100 stimuli — FEWER than HC', 'size': 11, 'color': GRAY},
            {'text': '• Fixations clustered in center, restricted area (convex hull small)', 'size': 11, 'color': GRAY},
            {'text': '• Per-category profile: BiCA-HS P(SZ) ≈ 0.80–1.00 uniformly high', 'size': 11, 'color': GRAY},
            {'text': '• GNN-CEFAM and ST-GNN also confident → all models agree', 'size': 11, 'color': GRAY},
            {'text': '• Final BiCA-HS overall: P(SZ) = 0.9990 (extreme SZ signal)', 'size': 11, 'bold': True, 'color': RED},
            {'text': '', 'size': 3},
            {'text': 'HC Subject #8 (RIGHT COLUMN):', 'size': 13, 'bold': True, 'color': BLUE},
            {'text': '• Scanpath: 1,360 fixations — MORE and more distributed', 'size': 11, 'color': GRAY},
            {'text': '• Fixations spread across screen (high convex hull area)', 'size': 11, 'color': GRAY},
            {'text': '• Per-category profile: P(SZ) ≈ 0.001–0.20 across all categories', 'size': 11, 'color': GRAY},
            {'text': '• Final BiCA-HS overall: P(SZ) = 0.0008 (extreme HC signal)', 'size': 11, 'bold': True, 'color': BLUE},
            {'text': '', 'size': 3},
            {'text': 'Population context (bottom panel):', 'size': 12, 'bold': True, 'color': PURP},
            {'text': '• All 160 subjects ranked by BiCA-HS P(SZ)', 'size': 11, 'color': GRAY},
            {'text': '• Clear bimodal separation: HC cluster near 0, SZ near 1', 'size': 11, 'color': GRAY},
            {'text': '• Ambiguous zone (P≈0.4–0.6): only ~8 subjects (5%)', 'size': 11, 'bold': True, 'color': GREEN},
        ])
add_pic(s, img("case_study_sz_vs_hc.png"), 0.1, 4.55, 9.8, 2.3)
footer(s, 45)
print("  Slide 45 (case study) added")

# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
print(f"\nSaving to {DST}...")
prs.save(DST)
total = len(prs.slides)
print(f"Done. Total slides: {total}")
print()
print("=" * 60)
print("SLIDE INVENTORY — ALL VISUALIZATIONS PLACED")
print("=" * 60)
print(f"  31 existing (updated) + 14 new = {total} slides")
print()
print("NEW SLIDES ADDED (32–45):")
print("  32: Tier 1 — 2D Scanpath Filtering (tier1_2d_filtering.png)")
print("  33: Tier 1 — Spatial Filter X/Y distributions")
print("  34: Tier 2B — Delta Feature Density Plots (6 KDE panels)")
print("  35: Tier 2 — Feature Violin Plots (8 features)")
print("  36: Tier 3 — XGBoost OOF Confusion Matrix")
print("  37: SHAP Waterfall — SZ#302 vs HC#8 individual")
print("  38: Feature Profile — z-score comparison SZ vs HC")
print("  39: Tier 4 — Threshold Sensitivity (ST-GNN vs BiCA-HS)")
print("  40: Tier 4 — Per-Category P(SZ) all 4 models")
print("  41: Tier 5 — 2D Ensemble Decision Space + LR boundary")
print("  42: Tier 5 — Comprehensive ROC + Reliability Diagram")
print("  43: OOF Decision Space Separation (class separation viz)")
print("  44: Ensemble Decision Space: CatBoost × BiCA-HS scatter")
print("  45: Case Study — SZ#302 vs HC#8 full panel")
print()
print("IMAGES ALREADY IN EXISTING SLIDES:")
print("  Slides 8,9: tier1_spatial/temporal filter figures")
print("  Slide 9: tier2_correlation_heatmap + tier3_shap_summary")
print("  Slide 10: tier2_delta_mechanism + tier2_feature_violin")
print("  Slide 11: tier3_shap_summary + tier3_tabular_comparison")
print("  Slides 13: tier4_graph_topology + tier4_scanpath_attention")
print("  Slide 15: tier5_reliability_diagram + tier5_dca")
print("  Slides 17: F7_roc_curves + tier5_roc_comparison")
print("  Slides 19–23: F2,F3,F5 ablation figures")
print("  Slide 21: tier3_shap_summary + explainability_shap_waterfall")
print("  Slide 22: tier4_real_attention + tier4_scanpath_attention")
print("  Slide 23: F5_fold_stability + tier5_multiseed_stability")
print("  Slide 28: tier4_resnet_tsne + tier4_resnet_umap (UPDATED)")
print("  Slide 29: tier1_pupil_dynamics")
print("  Slide 30: tier2_correlation_heatmap + tier2_feature_violin")
print()
print("IMAGES NOT USED (redundant/less critical):")
print("  tier3_feature_importance.png (covered by SHAP summary)")
print("  tier5_bica_s123/456 subfolders (covered by s42)")
print("  tier5_cefam_* (CEFAM results secondary)")
