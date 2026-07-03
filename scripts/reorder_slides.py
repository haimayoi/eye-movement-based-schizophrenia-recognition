"""
Reorder DAV_slides_v2.pptx into logical presentation flow
and fix all footer page numbers to X / 45.
"""
import sys, re
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

PPTX = r"D:\DAV\Eye Movement-Based Schizophrenia Recognition\DAV_slides_v2.pptx"
GRAY = RGBColor(0x40, 0x40, 0x40)

prs = Presentation(PPTX)
N = len(prs.slides)
print(f"Loaded {N} slides.")

# ── Print current titles for verification ─────────────────────────────
def get_title(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph.text.strip().replace('\n', ' ')[:70]
    return "(no placeholder title)"

print("\nCURRENT ORDER (0-based index → title):")
for i, s in enumerate(prs.slides):
    print(f"  [{i:2d}] {get_title(s)}")

# ══════════════════════════════════════════════════════════════════════
# NEW ORDER — 0-based indices of CURRENT slides → desired new position
# ══════════════════════════════════════════════════════════════════════
#
# Target logical structure (45 slides):
#  1  Title
#  2  Table of Contents
#  3  Motivation
#  4  Dataset & EMS
#  5  Related Works
#  6  Key Contributions
#  7  5-Tier Pipeline Overview
#
# ── TIER 1 ──
#  8  Tier 1 Preprocessing overview         (old 8  → idx 7)
#  9  Tier 1 2D Scanpath Filtering          (old 32 → idx 31)
# 10  Tier 1 Spatial Filter X/Y             (old 33 → idx 32)
# 11  Pupil Dynamics & Normalization        (old 29 → idx 28)
#
# ── TIER 2 ──
# 12  Tier 2 Feature Engineering overview   (old 9  → idx 8)
# 13  Feature Distributions & Correlation   (old 30 → idx 29)
# 14  Tier 2 Top-8 Violin Plots             (old 35 → idx 34)
# 15  Tier 2B Delta Features overview       (old 10 → idx 9)
# 16  Tier 2B Delta Density (6 KDE)         (old 34 → idx 33)
#
# ── TIER 3 ──
# 17  Tier 3 Tabular Models overview        (old 11 → idx 10)
# 18  Tier 3 OOF Confusion Matrix           (old 36 → idx 35)
# 19  OOF Decision Space Separation         (old 43 → idx 42)
#
# ── EVALUATION PROTOCOL ──
# 20  Evaluation Protocol                   (old 12 → idx 11)
#
# ── TIER 4 ──
# 21  Tier 4A GNN-CEFAM                     (old 13 → idx 12)
# 22  t-SNE + UMAP (motivates deep)         (old 28 → idx 27)
# 23  Tier 4B BiCA-HS                       (old 14 → idx 13)
# 24  Threshold Sensitivity (ST-GNN/BiCA)   (old 39 → idx 38)
# 25  Per-Category P(SZ) 4 models           (old 40 → idx 39)
#
# ── TIER 5 ──
# 26  Tier 5 Meta-Learner                   (old 15 → idx 14)
# 27  2D Ensemble Decision Space            (old 41 → idx 40)
#
# ── RESULTS ──
# 28  Main Results Table                    (old 16 → idx 15)
# 29  ROC Curves & DCA                      (old 17 → idx 16)
# 30  Comprehensive ROC + Calibration       (old 42 → idx 41)
# 31  Ablation Study                        (old 19 → idx 18)
# 32  Stimulus Category Analysis            (old 20 → idx 19)
# 33  Statistical Testing                   (old 18 → idx 17)
# 34  Multi-Seed Stability                  (old 23 → idx 22)
#
# ── XAI & CASE STUDY ──
# 35  Explainability SHAP Analysis          (old 21 → idx 20)
# 36  SHAP Waterfall (individual)           (old 37 → idx 36)
# 37  Explainability Attention Viz          (old 22 → idx 21)
# 38  Individual Feature Profile            (old 38 → idx 37)
# 39  Ensemble Decision Scatter             (old 44 → idx 43)
# 40  Case Study SZ#302 vs HC#8            (old 45 → idx 44)
#
# ── CURRICULUM & CONCLUSION ──
# 41  DAV Curriculum Coverage               (old 31 → idx 30)
# 42  Limitations                           (old 24 → idx 23)
# 43  Conclusion & Future Work              (old 25 → idx 24)
# 44  References                            (old 26 → idx 25)
# 45  Q&A                                  (old 27 → idx 26)

NEW_ORDER = [
    0,  # 1:  Title
    1,  # 2:  ToC
    2,  # 3:  Motivation
    3,  # 4:  Dataset
    4,  # 5:  Related Works
    5,  # 6:  Contributions
    6,  # 7:  Pipeline Overview
    # TIER 1
    7,  # 8:  Tier 1 Preprocessing
    31, # 9:  2D Filtering
    32, # 10: Spatial Filter X/Y
    28, # 11: Pupil Dynamics
    # TIER 2
    8,  # 12: Tier 2 Feature Eng
    29, # 13: Feature Dist & Corr
    34, # 14: Violin Plots (8 features)
    9,  # 15: Tier 2B Delta
    33, # 16: Delta Density (6 KDE)
    # TIER 3
    10, # 17: Tier 3 Tabular
    35, # 18: OOF Confusion Matrix
    42, # 19: OOF Decision Space
    # EVAL PROTOCOL
    11, # 20: Evaluation Protocol
    # TIER 4
    12, # 21: GNN-CEFAM
    27, # 22: t-SNE + UMAP
    13, # 23: BiCA-HS
    38, # 24: Threshold Sensitivity
    39, # 25: Per-Category P(SZ)
    # TIER 5
    14, # 26: Tier 5 Meta-Learner
    40, # 27: 2D Decision Space
    # RESULTS
    15, # 28: Main Results Table
    16, # 29: ROC + DCA
    41, # 30: Comprehensive ROC + Calibration
    18, # 31: Ablation Study
    19, # 32: Category Analysis
    17, # 33: Statistical Testing
    22, # 34: Multi-Seed Stability
    # XAI & CASE STUDY
    20, # 35: SHAP Analysis
    36, # 36: SHAP Waterfall
    21, # 37: Attention Viz
    37, # 38: Feature Profile
    43, # 39: Ensemble Decision Scatter
    44, # 40: Case Study
    # CURRICULUM & CONCLUSION
    30, # 41: DAV Curriculum
    23, # 42: Limitations
    24, # 43: Conclusion
    25, # 44: References
    26, # 45: Q&A
]

assert len(NEW_ORDER) == N, f"Order list has {len(NEW_ORDER)} items, expected {N}"
assert sorted(NEW_ORDER) == list(range(N)), "ORDER must be a permutation of 0..N-1"

# ── Reorder the XML slide ID list ──────────────────────────────────────
sldIdLst = prs.slides._sldIdLst
all_sldId = list(sldIdLst)          # snapshot of <p:sldId> elements

# Remove all
for elem in all_sldId:
    sldIdLst.remove(elem)

# Re-append in new order
for idx in NEW_ORDER:
    sldIdLst.append(all_sldId[idx])

print("\nNEW ORDER after reordering:")
for i, s in enumerate(prs.slides, 1):
    print(f"  [{i:2d}] {get_title(s)}")

# ══════════════════════════════════════════════════════════════════════
# FIX FOOTER PAGE NUMBERS — scan every shape on every slide
# Look for text matching  \d+\s*/\s*\d+  in textboxes near bottom
# and replace with   new_slide_num / 45
# ══════════════════════════════════════════════════════════════════════
TOTAL = 45
PAGE_PATTERN = re.compile(r'^\d+\s*/\s*\d+$')

fixed = 0
for slide_num, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            full_text = ''.join(r.text for r in para.runs).strip()
            if PAGE_PATTERN.match(full_text):
                # Replace ALL runs with single corrected run
                new_text = f'{slide_num} / {TOTAL}'
                # Clear existing runs
                for run in para.runs:
                    run.text = ''
                # Set text in first run
                if para.runs:
                    para.runs[0].text = new_text
                    para.runs[0].font.size = Pt(9)
                    para.runs[0].font.color.rgb = GRAY
                fixed += 1
                print(f"  Slide {slide_num:2d}: footer fixed → '{new_text}'")

print(f"\nFixed {fixed} footer numbers.")

# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
prs.save(PPTX)
print(f"\nSaved {TOTAL} slides to {PPTX}")
print("\nFINAL SLIDE ORDER:")
for i, s in enumerate(prs.slides, 1):
    print(f"  {i:2d}. {get_title(s)}")
