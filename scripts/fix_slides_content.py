"""
Targeted fix for DAV_slides_v2.pptx:
1. Fix content inaccuracies vs actual results
2. Move Evaluation Protocol to slide 8 (after Pipeline Overview)
3. Remove duplicate textboxes on slide 33
4. Re-fix all footer page numbers
"""
import sys, re
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

PPTX = r"D:\DAV\Eye Movement-Based Schizophrenia Recognition\DAV_slides_v2.pptx"
prs = Presentation(PPTX)
N = len(prs.slides)
print(f"Loaded {N} slides.")

RED   = RGBColor(0xC0, 0x00, 0x00)
BLUE  = RGBColor(0x1F, 0x49, 0x7D)
GRAY  = RGBColor(0x40, 0x40, 0x40)
GREEN = RGBColor(0x00, 0x70, 0x00)
ORANGE= RGBColor(0xC0, 0x50, 0x00)

def get_title(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph.text.strip()[:70]
    return "(no title)"

def rebuild_tf(shape, lines, size=12, color=GRAY):
    """Rebuild entire text frame with given lines list. Each line: str or dict."""
    tf = shape.text_frame
    tf.word_wrap = True
    # Remove all paragraphs except the first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()

    from pptx.enum.text import PP_ALIGN
    first = True
    for line in lines:
        if isinstance(line, dict):
            text  = line.get('text', '')
            sz    = line.get('size', size)
            bold  = line.get('bold', False)
            col   = line.get('color', color)
            itl   = line.get('italic', False)
            space = line.get('space', 0)
        else:
            text, sz, bold, col, itl, space = line, size, False, color, False, 0

        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.clear()
        if space:
            para.space_before = Pt(space)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.italic = itl

# ══════════════════════════════════════════════════════════════════════
# FIX 1 — Slide 17: Tier 3 Tabular — add LightGBM AUC
# ══════════════════════════════════════════════════════════════════════
s17 = prs.slides[16]  # 0-based
print(f"\nFixing Slide 17: {get_title(s17)}")
for shape in s17.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if '135-dim expert vector' in txt:
            rebuild_tf(shape, [
                {'text': 'INPUT: 135-dim expert vector (15 stim-level + 120 cat/delta features)', 'size': 12, 'bold': True, 'color': RED},
                {'text': '', 'size': 4},
                {'text': 'Three GBDT baselines trained on subject-level aggregated features:', 'size': 12, 'bold': True, 'color': BLUE},
                {'text': '', 'size': 3},
                {'text': '• XGBoost:    AUC = 0.872 ± 0.005  |  ACC = 77.7 ± 1.1%', 'size': 12, 'color': GRAY},
                {'text': '• LightGBM:   AUC = 0.843 ± 0.003  |  ACC = 77.7 ± 1.1%', 'size': 12, 'color': GRAY},
                {'text': '• CatBoost:   AUC = 0.899 ± 0.001  |  ACC = 81.0 ± 0.8%  ← best', 'size': 12, 'bold': True, 'color': BLUE},
                {'text': '', 'size': 3},
                {'text': 'Protocol: 4-fold GroupKFold (subject-level, zero data leakage)', 'size': 11, 'color': GRAY},
                {'text': 'Aggregation: mean of 4 per-category probability scores → subject P(SZ)', 'size': 11, 'color': GRAY},
                {'text': '', 'size': 4},
                {'text': 'Feature Importance (DAV: C7 — XAI / SHAP)', 'size': 12, 'bold': True, 'color': RED},
                {'text': '• SHAP beeswarm: Synthetic_pupil_std #1 most important', 'size': 11, 'color': GRAY},
                {'text': '• Pupil variability features dominate (7 of top 10)', 'size': 11, 'color': GRAY},
                {'text': '• Delta(Soc−Nat) features in top 10 → validates Tier 2B design', 'size': 11, 'color': GRAY},
                {'text': '', 'size': 4},
                {'text': 'OUTPUT: OOF subject P_tab → Tier 5 meta-learner', 'size': 12, 'bold': True, 'color': GREEN},
                {'text': 'CatBoost selected as Tier 5 input (best AUC)', 'size': 11, 'color': GRAY},
            ])
            print("  → Left content rebuilt with LightGBM AUC")
            break

# ══════════════════════════════════════════════════════════════════════
# FIX 2 — Slide 31: Ablation — fix truncated text & accurate narrative
# ══════════════════════════════════════════════════════════════════════
s31 = prs.slides[30]
print(f"\nFixing Slide 31: {get_title(s31)}")
for shape in s31.shapes:
    if shape.has_text_frame and 'Progressive Architecture Ablation' in shape.text_frame.text:
        rebuild_tf(shape, [
            {'text': 'Progressive Architecture Ablation (seed 42)', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 4},
            {'text': 'Stage 1 — XGBoost (tabular baseline):', 'size': 12, 'bold': True, 'color': GRAY},
            {'text': '  AUC = 0.870 | No temporal or graph modeling', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 2},
            {'text': 'Stage 2 — CatBoost (best tabular):', 'size': 12, 'bold': True, 'color': GRAY},
            {'text': '  AUC = 0.899  Δ+0.029 | Gradient-boosted trees with categorical features', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 2},
            {'text': 'Stage 3 — ST-GNN (graph topology only, no expert stream):', 'size': 12, 'bold': True, 'color': GRAY},
            {'text': '  AUC = 0.933  Δ+0.034 | Spatiotemporal graph adds significant value', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 2},
            {'text': 'Stage 4 — GNN-CEFAM (GAT + CEFAM cross-attention + expert):', 'size': 12, 'bold': True, 'color': ORANGE},
            {'text': '  AUC = 0.921  Δ−0.012 vs ST-GNN', 'size': 11, 'color': ORANGE},
            {'text': '  Direction 2 degeneracy (seq_len=1) neutralises CEFAM benefit', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 2},
            {'text': 'Stage 5 — BiCA-HS (Transformer + bidirectional cross-attention):', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': '  AUC = 0.955  Δ+0.022 vs ST-GNN  |  Δ+0.034 vs GNN-CEFAM', 'size': 11, 'bold': True, 'color': BLUE},
            {'text': '  Processes 200 fixations; both attention directions non-trivial', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 2},
            {'text': 'Stage 6 — Tier 5 Ensemble (CatBoost + BiCA-HS):', 'size': 12, 'bold': True, 'color': GREEN},
            {'text': '  AUC = 0.959  Δ+0.004 | Not statistically significant (p=0.54)', 'size': 11, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'BiCA-HS is the optimal single model: AUC=0.955, Sens=87.5%, Spec=87.1%', 'size': 12, 'bold': True, 'color': GREEN},
        ])
        print("  → Ablation text rebuilt with accurate progression")
        break

# ══════════════════════════════════════════════════════════════════════
# FIX 3 — Slide 32: Category Analysis — fix wrong AUC values
# Actual F3 data (seed 42): Social=0.958, Manip=0.956, Natural=0.956, Synth=0.958
# ══════════════════════════════════════════════════════════════════════
s32 = prs.slides[31]
print(f"\nFixing Slide 32: {get_title(s32)}")
for shape in s32.shapes:
    if shape.has_text_frame and 'Per-Category OOF AUC' in shape.text_frame.text:
        rebuild_tf(shape, [
            {'text': 'Per-Category OOF AUC — BiCA-HS (seed 42, F3 ablation data)', 'size': 13, 'bold': True, 'color': RED},
            {'text': '', 'size': 4},
            {'text': 'Social   (portraits, people):        AUC = 0.958', 'size': 12, 'color': GRAY},
            {'text': 'Manipulated (edited scenes):         AUC = 0.956', 'size': 12, 'color': GRAY},
            {'text': 'Natural  (landscapes):               AUC = 0.956', 'size': 12, 'color': GRAY},
            {'text': 'Synthetic (computer graphics):       AUC = 0.958', 'size': 12, 'color': GRAY},
            {'text': '', 'size': 5},
            {'text': 'Key finding: BiCA-HS is CATEGORY-ROBUST', 'size': 13, 'bold': True, 'color': GREEN},
            {'text': 'AUC variation across categories < 0.002 → near-zero category bias', 'size': 12, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'Tabular models (XGBoost/CatBoost):', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': 'Identical AUC across all 4 categories (subject-level aggregation)', 'size': 11, 'color': GRAY},
            {'text': 'XGBoost: 0.870 | CatBoost: 0.898 uniformly', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'ST-GNN category profile (seed 42):', 'size': 12, 'bold': True, 'color': BLUE},
            {'text': 'Social=0.916, Manip=0.933, Natural=0.940, Synth=0.946', 'size': 11, 'color': GRAY},
            {'text': 'Natural/Synthetic slightly better → non-social scenes easier for GNN', 'size': 11, 'italic': True, 'color': GRAY},
            {'text': '', 'size': 4},
            {'text': 'Validates Tier 2B: category-specific delta features', 'size': 12, 'bold': True, 'color': GREEN},
            {'text': 'carry the discriminative signal — robust across all stimulus types', 'size': 11, 'color': GRAY},
        ])
        print("  → Category analysis AUC values corrected (0.958/0.956/0.956/0.958)")
        break

# ══════════════════════════════════════════════════════════════════════
# FIX 4 — Slide 33: Statistical Testing — remove duplicate textboxes
# Shapes[5] (duplicate DeLong) and [7] (duplicate multi-seed) should go
# ══════════════════════════════════════════════════════════════════════
s33 = prs.slides[32]
print(f"\nFixing Slide 33: {get_title(s33)}")

# Collect textboxes (type 17) by content, find duplicates
seen_texts = {}
to_remove = []
for shape in s33.shapes:
    if shape.shape_type == 17 and shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if txt in seen_texts:
            to_remove.append(shape)
            print(f"  → Removing duplicate textbox: {repr(txt[:60])}")
        else:
            seen_texts[txt] = shape

for shape in to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

print(f"  Removed {len(to_remove)} duplicate shapes")

# ══════════════════════════════════════════════════════════════════════
# FIX 5 — Slide 28: Main Results — verify table is complete
# The table text was truncated in print. Check if slide 28 content is complete.
# ══════════════════════════════════════════════════════════════════════
s28 = prs.slides[27]
print(f"\nChecking Slide 28: {get_title(s28)}")
for shape in s28.shapes:
    if shape.has_text_frame and 'TABLE 1' in shape.text_frame.text:
        full = shape.text_frame.text
        if 'Tier5' not in full or 'BiCA-HS' not in full:
            print("  WARNING: Table might be incomplete!")
        else:
            print(f"  Table OK, length={len(full)} chars")
        # Show last 200 chars
        print(f"  ...{repr(full[-300:])}")

# ══════════════════════════════════════════════════════════════════════
# REORDER — Move Evaluation Protocol from slide 20 (idx 19) to slide 8 (idx 7)
# Current 0-based: 0..6 = Title..Pipeline | 7=Tier1 | ... | 19=EvalProt | 20=CEFAM
# New 0-based:     0..6 = Title..Pipeline | 7=EvalProt | 8=Tier1 | ... | 19=CEFAM
# ══════════════════════════════════════════════════════════════════════
print(f"\n--- REORDERING: Moving Eval Protocol to slide 8 ---")
print(f"  Current slide 20: {get_title(prs.slides[19])}")

# New order: indices 0-6, then 19, then 7-18, then 20-44
NEW_ORDER = list(range(7)) + [19] + list(range(7, 19)) + list(range(20, 45))
assert len(NEW_ORDER) == 45
assert sorted(NEW_ORDER) == list(range(45))

sldIdLst = prs.slides._sldIdLst
all_sldId = list(sldIdLst)
for elem in all_sldId:
    sldIdLst.remove(elem)
for idx in NEW_ORDER:
    sldIdLst.append(all_sldId[idx])

print(f"  New slide 8: {get_title(prs.slides[7])}")
print(f"  New slide 9: {get_title(prs.slides[8])}")

# ══════════════════════════════════════════════════════════════════════
# FIX 6 — Update all footer page numbers after reorder
# ══════════════════════════════════════════════════════════════════════
TOTAL = 45
PAGE_PATTERN = re.compile(r'^\d+\s*/\s*\d+$')
fixed = 0
for slide_num, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = ''.join(r.text for r in para.runs).strip()
            if PAGE_PATTERN.match(full):
                new_text = f'{slide_num} / {TOTAL}'
                for r in para.runs:
                    r.text = ''
                if para.runs:
                    para.runs[0].text = new_text
                    para.runs[0].font.size = Pt(9)
                    para.runs[0].font.color.rgb = GRAY
                fixed += 1

print(f"\nFixed {fixed} footer numbers (1/{TOTAL} ... {TOTAL}/{TOTAL})")

# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
out = PPTX.replace('.pptx', '_fixed.pptx')
prs.save(out)
print(f"\nSaved to {out}")
print(f"\nFINAL SLIDE ORDER:")
for i, s in enumerate(prs.slides, 1):
    print(f"  {i:2d}. {get_title(s)}")
